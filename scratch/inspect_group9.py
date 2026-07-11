from pptx import Presentation

path = "C:/Users/ROHITH P/Downloads/Internship Certificate.pptx"
prs = Presentation(path)
slide = prs.slides[0]

group_shape = None
for shape in slide.shapes:
    if shape.shape_id == 9:
        group_shape = shape
        break

if group_shape is None:
    print("Group shape ID 9 not found")
else:
    print(f"Group 9 Bounding Box: L={group_shape.left/914400:.3f}, T={group_shape.top/914400:.3f}, W={group_shape.width/914400:.3f}, H={group_shape.height/914400:.3f}")
    if hasattr(group_shape, 'shapes'):
        print(f"Sub-shapes inside Group 9:")
        min_left = float('inf')
        max_right = float('-inf')
        min_top = float('inf')
        max_bottom = float('-inf')
        
        for j, s in enumerate(group_shape.shapes):
            l = s.left / 914400
            t = s.top / 914400
            w = s.width / 914400
            h = s.height / 914400
            r = l + w
            b = t + h
            
            min_left = min(min_left, l)
            max_right = max(max_right, r)
            min_top = min(min_top, t)
            max_bottom = max(max_bottom, b)
            
            print(f"  Sub {j}: {s.name} | L={l:.3f}, T={t:.3f}, W={w:.3f}, H={h:.3f}")
            
        print(f"\nOpaque Bounds calculated from sub-shapes:")
        print(f"  Left  : {min_left:.3f} in")
        print(f"  Right : {max_right:.3f} in")
        print(f"  Top   : {min_top:.3f} in")
        print(f"  Bottom: {max_bottom:.3f} in")
    else:
        print("Shape is not a group shape")
