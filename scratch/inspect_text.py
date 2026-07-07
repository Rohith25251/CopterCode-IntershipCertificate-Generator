import os
import glob
from pptx import Presentation

# Find any PPTX files in Downloads or Desktop or current directory
paths = [
    "C:/Users/ROHITH P/Downloads/*.pptx",
    "C:/Users/ROHITH P/Desktop/*.pptx",
    "./*.pptx",
    "../**/*.pptx"
]

found = []
for p in paths:
    found.extend(glob.glob(p, recursive=True))

print("Found PPTX files:")
for f in found:
    print(f)
    try:
        prs = Presentation(f)
        slide = prs.slides[0]
        print(f"  Successfully opened {f}")
        for i, s in enumerate(slide.shapes):
            if s.has_text_frame:
                for j, p in enumerate(s.text_frame.paragraphs):
                    text = p.text
                    if "notify" in text or "successful" in text or "DEPARTMENT" in text:
                        print(f"  Shape {i}, Paragraph {j}: {repr(text)}")
                        for k, r in enumerate(p.runs):
                            print(f"    Run {k} (bold={r.font.bold}, size={r.font.size}): {repr(r.text)}")
    except Exception as e:
        print(f"  Error reading {f}: {e}")
