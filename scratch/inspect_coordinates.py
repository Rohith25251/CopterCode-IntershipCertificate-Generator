import os
from pptx import Presentation

templates = {
    "Internship": "C:/Users/ROHITH P/Downloads/Internship Certificate.pptx",
    "Experience": "C:/Users/ROHITH P/Downloads/Experience Certificate.pptx",
    "LOR": "C:/Users/ROHITH P/Downloads/Letter Of Recommandation.pptx"
}

for name, path in templates.items():
    if not os.path.exists(path):
        print(f"File not found: {path}")
        continue
        
    print(f"\n==========================================")
    print(f"TEMPLATES: {name} ({path})")
    print(f"==========================================")
    
    prs = Presentation(path)
    slide = prs.slides[0]
    
    # Print slide dimensions
    slide_w_in = prs.slide_width / 914400
    slide_h_in = prs.slide_height / 914400
    print(f"Slide Dimensions: {slide_w_in:.3f} x {slide_h_in:.3f} inches")
    
    for i, shape in enumerate(slide.shapes):
        shape_id = shape.shape_id
        left_in = shape.left / 914400
        top_in = shape.top / 914400
        width_in = shape.width / 914400
        height_in = shape.height / 914400
        
        has_text = "YES" if shape.has_text_frame else "NO"
        text_snippet = ""
        if shape.has_text_frame:
            text_snippet = shape.text_frame.text.replace('\n', '\\n')[:60]
            
        print(f"Idx: {i:2d} | ID: {shape_id:3d} | Name: {shape.name:25s} | HasText: {has_text} | L: {left_in:6.3f} | T: {top_in:6.3f} | W: {width_in:6.3f} | H: {height_in:6.3f} | Text: {text_snippet}")
