import os
import re
from pptx import Presentation

template_path = "C:/Users/ROHITH P/Downloads/certificate template.pptx"
out_pptx = "scratch/test_fixed.pptx"
out_pdf = "scratch/test_fixed.pdf"

prs = Presentation(template_path)
slide = prs.slides[0]

replacements = {
    "<<NAME>>": "GHAANDHAMADHAN V",
    "<<INSTITUTION>>": "KONGU ENGINEERING COLLEGE",
    "<<COLLEGE>>": "KONGU ENGINEERING COLLEGE",
    "<<YEAR>>": "3rd Year",
    "<<DEPARTMENT>>": "COMPUTER SCIENCE AND ENGINEERING",
    "<<DOMAIN>>": "FULL-STACK Development & Software Development",
    "<<ROLE>>": "FULL-STACK Development & Software Development",
    "<<PROJECT>>": "FULL-STACK DEVELOPMENT & SOFTWARE DEVELOPMENT - ERP SOLUTION & SITES",
    "<<INTERNSHIP & LIVE PROJECT AREA>>": "FULL-STACK DEVELOPMENT & SOFTWARE DEVELOPMENT - ERP SOLUTION & SITES",
    "<<BATCH>>": "1 MONTH - JAN",
    "<<BATCH >>": "1 MONTH - JAN",
    "<<DATE>>": "08.07.2026",
    "<<DT>>": "08.07.2026"
}

for shape in list(slide.shapes):
    if shape.has_text_frame:
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # Replace placeholders
                for key, val in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, val)
                    else:
                        guill_key = key.replace("<<", "«").replace(">>", "»")
                        if guill_key in run.text:
                            run.text = run.text.replace(guill_key, val)
                
                # Add space after comma if missing (e.g. ",COMPUTER" -> ", COMPUTER")
                run.text = re.sub(r',([a-zA-Z])', r', \1', run.text)
                
                # Clean up multiple spaces (replace 2 or more spaces/tabs/newlines with a single space)
                run.text = re.sub(r'\s{2,}', ' ', run.text)

prs.save(out_pptx)

import win32com.client
powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
presentation = powerpoint.Presentations.Open(os.path.abspath(out_pptx), WithWindow=False)
presentation.SaveAs(os.path.abspath(out_pdf), 32)
presentation.Close()
powerpoint.Quit()

import fitz
doc = fitz.open(out_pdf)
for page in doc:
    print("--- PDF Page Text ---")
    print(page.get_text())
doc.close()
