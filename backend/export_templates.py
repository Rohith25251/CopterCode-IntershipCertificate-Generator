import os
import re
import json
import uuid
import shutil
import subprocess
from pptx import Presentation

# Paths to input templates
TEMPLATES = {
    "internship": "C:/Users/ROHITH P/Downloads/Internship Certificate.pptx",
    "experience": "C:/Users/ROHITH P/Downloads/Experience Certificate.pptx",
    "lor": "C:/Users/ROHITH P/Downloads/Letter Of Recommandation.pptx"
}

OUTPUT_DIR = "backend/templates"

def has_placeholder(text):
    return any(p in text for p in ["<<", ">>", "«", "»"])

def is_body_text_shape(shape):
    if not shape.has_text_frame:
        return False
    # Exclude elements at the very bottom edge (footer area, T > 11.0) from flowing
    if shape.top.inches > 11.0:
        return False
    text = shape.text_frame.text.strip()
    if not text:
        return False
        
    # Always include if it contains a placeholder
    if has_placeholder(text):
        return True
        
    # Bounding check for vertical column shapes (wide text blocks in the middle of A4)
    left_in = shape.left.inches
    top_in = shape.top.inches
    width_in = shape.width.inches
    
    is_wide_and_centered = (left_in < 1.0) and (width_in > 6.5)
    is_in_middle_band = (top_in > 2.2) and (top_in < 9.0)
    
    return is_wide_and_centered and is_in_middle_band

def get_alignment_str(paragraph):
    from pptx.enum.text import PP_ALIGN
    align = paragraph.alignment
    if align == PP_ALIGN.CENTER:
        return "center"
    elif align == PP_ALIGN.RIGHT:
        return "right"
    elif align == PP_ALIGN.JUSTIFY:
        return "justify"
    return "left"

def get_color_hex(run):
    try:
        color = run.font.color
        if color and color.type == 1:  # RGBColor
            rgb = color.rgb
            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
    except:
        pass
    return None

_EMU_PER_INCH = 914400

def get_shape_fill_color(shape):
    """Extract the dominant fill color from a shape (including groups)."""
    try:
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        for srgb in shape._element.iter(f'{{{a_ns}}}srgbClr'):
            val = (srgb.get('val') or '').lower()
            if val and val != 'ffffff':
                return f'#{val}'
    except Exception:
        pass
    return '#000000'

def extract_group_text_children(group_shape):
    """Extract text shapes from inside a GROUP shape with absolute slide coords."""
    try:
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        el = group_shape._element
        grp_sp_pr = next((c for c in el if c.tag.endswith('}grpSpPr')), None)
        if grp_sp_pr is None: return []
        xfrm = next((c for c in grp_sp_pr if c.tag.endswith('}xfrm')), None)
        if xfrm is None: return []
        ns = f'{{{a_ns}}}'
        def _i(elem, attr, d=0): return int(elem.get(attr, d))
        off = xfrm.find(f'{ns}off'); ext = xfrm.find(f'{ns}ext')
        chOff = xfrm.find(f'{ns}chOff'); chExt = xfrm.find(f'{ns}chExt')
        if any(x is None for x in [off, ext, chOff, chExt]): return []
        gx = _i(off,'x'); gy = _i(off,'y'); gcx = _i(ext,'cx'); gcy = _i(ext,'cy')
        cox = _i(chOff,'x'); coy = _i(chOff,'y')
        ccx = _i(chExt,'cx',1) or 1; ccy = _i(chExt,'cy',1) or 1
        sx = gcx/ccx; sy = gcy/ccy
        try: child_shapes = group_shape.shapes
        except Exception: return []
        result = []
        for child in child_shapes:
            if not child.has_text_frame: continue
            text = child.text_frame.text.strip()
            if not text: continue
            left_in = (gx + (child.left - cox) * sx) / _EMU_PER_INCH
            top_in  = (gy + (child.top  - coy) * sy) / _EMU_PER_INCH
            width_in = max(child.width * sx / _EMU_PER_INCH, 0.5)
            height_in = max(child.height * sy / _EMU_PER_INCH, 0.2)
            font_name = "Calibri"; font_size = 14
            font_color = "#000000"; bold = False; italic = False; align = "left"
            if child.text_frame.paragraphs:
                p0 = child.text_frame.paragraphs[0]
                align = get_alignment_str(p0)
                if p0.runs:
                    r0 = p0.runs[0]
                    if r0.font.name: font_name = r0.font.name
                    if r0.font.size: font_size = r0.font.size.pt
                    c = get_color_hex(r0)
                    if c: font_color = c
                    bold = bool(r0.font.bold); italic = bool(r0.font.italic)
            paragraphs_cfg = []
            for p in child.text_frame.paragraphs:
                runs_cfg = []
                for r in p.runs:
                    r_color = get_color_hex(r) or font_color
                    runs_cfg.append({
                        "text": r.text,
                        "font_name": r.font.name or font_name,
                        "font_size": r.font.size.pt if r.font.size else font_size,
                        "bold": bool(r.font.bold), "italic": bool(r.font.italic),
                        "underline": bool(r.font.underline), "color": r_color
                    })
                paragraphs_cfg.append({"align": get_alignment_str(p), "runs": runs_cfg})
            result.append({
                "id": child.shape_id if hasattr(child, 'shape_id') else 0,
                "name": f"{group_shape.name}::{child.name}",
                "left": left_in, "top": top_in,
                "width": width_in, "height": height_in,
                "font_name": font_name, "font_size": font_size,
                "color": font_color, "bold": bold, "italic": italic,
                "align": align, "original_text": text,
                "paragraphs": paragraphs_cfg, "is_flow": (top_in <= 11.0)
            })
        return result
    except Exception as e:
        print(f"[Warning] extract_group_text_children failed for {group_shape.name}: {e}")
        return []

def process_template(template_name, pptx_path):
    print(f"\nProcessing template: {template_name} from {pptx_path}...")
    if not os.path.exists(pptx_path):
        print(f"Error: PPTX file not found at {pptx_path}")
        return
        
    os.makedirs(f"{OUTPUT_DIR}/{template_name}", exist_ok=True)
    
    # 1. Load original presentation
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    
    # Slide dimensions in inches
    slide_w_in = prs.slide_width.inches
    slide_h_in = prs.slide_height.inches
    
    layout_data = {
        "template": template_name,
        "width_in": slide_w_in,
        "height_in": slide_h_in,
        "shapes": []
    }
    
    # Track which shapes we need to clear from the background
    shapes_to_clear = []
    
    for i, shape in enumerate(slide.shapes):
        h_in = shape.height.inches if hasattr(shape, "height") else 0
        w_in = shape.width.inches if hasattr(shape, "width") else 0
        l_in = shape.left.inches if hasattr(shape, "left") else 0
        t_in = shape.top.inches if hasattr(shape, "top") else 0
        if not shape.has_text_frame:
            if h_in < 0.25 and w_in > 1.0 and l_in >= 0:
                line_color = get_shape_fill_color(shape)
                shape_cfg = {
                    "id": shape.shape_id if hasattr(shape, "shape_id") else len(layout_data["shapes"]),
                    "name": shape.name,
                    "left": l_in, "top": t_in, "width": w_in, "height": h_in,
                    "is_line": True, "is_qr": False,
                    "is_flow": t_in <= 11.0, "line_color": line_color
                }
                layout_data["shapes"].append(shape_cfg)
                shapes_to_clear.append(shape)
            elif shape.shape_type == 6 and 0.15 <= h_in <= 1.5 and 0.5 <= t_in <= 11.0:
                group_children = extract_group_text_children(shape)
                for child_cfg in group_children:
                    layout_data["shapes"].append(child_cfg)
                if group_children:
                    shapes_to_clear.append(shape)
            continue
            
        text = shape.text_frame.text.strip()
        if not text:
            continue
            
        print(f"  Exporting text shape {shape.shape_id} ({shape.name}) with text: {repr(text[:50])}...")
        
        # Get geometry
        left_in = shape.left.inches
        top_in = shape.top.inches
        width_in = shape.width.inches
        height_in = shape.height.inches
        
        # Default text styling from first paragraph/run
        font_name = "Calibri"
        font_size = 14
        font_color = "#000000"
        bold = False
        italic = False
        align = "left"
        
        if shape.text_frame.paragraphs:
            p = shape.text_frame.paragraphs[0]
            align = get_alignment_str(p)
            if p.runs:
                r = p.runs[0]
                if r.font.name:
                    font_name = r.font.name
                if r.font.size:
                    font_size = r.font.size.pt
                c = get_color_hex(r)
                if c:
                    font_color = c
                bold = bool(r.font.bold)
                italic = bool(r.font.italic)
                
        # Parse paragraphs and runs configuration to recreate the DOM layout accurately
        paragraphs_cfg = []
        for p in shape.text_frame.paragraphs:
            runs_cfg = []
            for r in p.runs:
                r_color = get_color_hex(r) or font_color
                runs_cfg.append({
                    "text": r.text,
                    "font_name": r.font.name or font_name,
                    "font_size": r.font.size.pt if r.font.size else font_size,
                    "bold": bool(r.font.bold),
                    "italic": bool(r.font.italic),
                    "underline": bool(r.font.underline),
                    "color": r_color
                })
            paragraphs_cfg.append({
                "align": get_alignment_str(p),
                "runs": runs_cfg
            })
            
        is_flow = is_body_text_shape(shape)
        shape_cfg = {
            "id": shape.shape_id,
            "name": shape.name,
            "left": left_in,
            "top": top_in,
            "width": width_in,
            "height": height_in,
            "font_name": font_name,
            "font_size": font_size,
            "color": font_color,
            "bold": bold,
            "italic": italic,
            "align": align,
            "original_text": text,
            "paragraphs": paragraphs_cfg,
            "is_flow": is_flow
        }
        
        layout_data["shapes"].append(shape_cfg)
        shapes_to_clear.append(shape)
        
    # Write layout.json
    json_path = f"{OUTPUT_DIR}/{template_name}/layout.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(layout_data, f, indent=2)
    print(f"  Wrote shape geometries to {json_path}")
    
    # 2. Clear placeholder text boxes and move line shapes off-slide
    from pptx.util import Inches
    for shape in shapes_to_clear:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.text = ""
        else:
            try:
                shape.left = Inches(-20)
            except Exception:
                pass
                
    temp_pptx = f"{OUTPUT_DIR}/{template_name}/temp_bg.pptx"
    prs.save(temp_pptx)
    
    # 3. Export to PNG
    png_path = f"{OUTPUT_DIR}/{template_name}/background.png"
    exported = False
    
    # Method A: Try PowerPoint COM on Windows
    try:
        import win32com.client
        print(f"  [Windows] Exporting slide directly to PNG via PowerPoint COM automation...")
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(os.path.abspath(temp_pptx), WithWindow=False)
        presentation.Slides(1).Export(os.path.abspath(png_path), "PNG")
        presentation.Close()
        powerpoint.Quit()
        exported = True
        print(f"  [Windows] Background image saved to {png_path}")
    except Exception as e:
        print(f"  [Windows] COM automation not available or failed: {e}")
        
    # Method B: Fallback to LibreOffice + PyMuPDF
    if not exported:
        temp_pdf_dir = f"{OUTPUT_DIR}/{template_name}"
        cmd = [
            "soffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", temp_pdf_dir,
            temp_pptx
        ]
        print(f"  [Linux/Docker] Converting PPTX to background PDF via LibreOffice...")
        try:
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            temp_pdf = f"{OUTPUT_DIR}/{template_name}/temp_bg.pdf"
            
            import fitz  # PyMuPDF
            print(f"  [Linux/Docker] Rendering background PDF to high-res PNG...")
            doc = fitz.open(temp_pdf)
            page = doc[0]
            matrix = fitz.Matrix(300 / 72, 300 / 72)  # 300 DPI
            pix = page.get_pixmap(matrix=matrix)
            pix.save(png_path)
            doc.close()
            
            if os.path.exists(temp_pdf):
                os.remove(temp_pdf)
            exported = True
            print(f"  [Linux/Docker] Background image saved to {png_path}")
        except Exception as e:
            print(f"  [Linux/Docker] LibreOffice rendering failed: {e}")
            
    # Clean up temporary presentation
    if os.path.exists(temp_pptx):
        os.remove(temp_pptx)
        
    if exported:
        print(f"Successfully processed {template_name}!")
    else:
        print(f"Failed to export background image for {template_name}")

if __name__ == "__main__":
    for name, path in TEMPLATES.items():
        process_template(name, path)
