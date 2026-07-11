import io
import os
import uuid
import re
import urllib.request
import tempfile
import pandas as pd
from typing import Dict, List, Optional
from datetime import datetime
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, status, Response, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from dotenv import load_dotenv
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

# win32com client has been removed to run reliably on headless Linux
import fitz
from pptx import Presentation
import asyncio
import json
try:
    from layout_engine import LayoutEngine
except ImportError:
    from backend.layout_engine import LayoutEngine

# QR generation library
import qrcode

# Supabase python client
from supabase import create_client, Client

# ReportLab & PyPDF Imports
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.styles import ParagraphStyle
from reportlab.platypus import Paragraph, Frame
from reportlab.platypus.flowables import KeepInFrame
from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
from reportlab.lib.utils import ImageReader
from pypdf import PdfReader, PdfWriter

load_dotenv()

def register_fonts():
    fonts_dir = os.path.join(os.path.dirname(__file__), "fonts")
    if not os.path.exists(fonts_dir):
        print("WARNING: Fonts directory not found at", fonts_dir)
        return
    for filename in os.listdir(fonts_dir):
        if filename.lower().endswith(".ttf"):
            font_name = os.path.splitext(filename)[0]
            font_path = os.path.join(fonts_dir, filename)
            try:
                pdfmetrics.registerFont(TTFont(font_name, font_path))
                print(f"Registered font: {font_name} from {font_path}")
            except Exception as e:
                print(f"Failed to register font {font_name}: {e}")

# Register fonts on startup
register_fonts()

import asyncio
import time
import tempfile

class SimpleTTLCache:
    def __init__(self, maxsize=200, ttl=600):
        self.maxsize = maxsize
        self.ttl = ttl
        self.cache = {}
        
    def get(self, key):
        if key in self.cache:
            val, expiry = self.cache[key]
            if time.time() < expiry:
                return val
            del self.cache[key]
        return None
        
    def set(self, key, val):
        if len(self.cache) >= self.maxsize:
            now = time.time()
            expired = [k for k, (v, exp) in self.cache.items() if now >= exp]
            if expired:
                for k in expired:
                    del self.cache[k]
            else:
                first_key = next(iter(self.cache))
                del self.cache[first_key]
        self.cache[key] = (val, time.time() + self.ttl)

template_cache = {}  # slot_id (e.g. "slot-1") -> pptx_bytes
pdf_cache = SimpleTTLCache(maxsize=200, ttl=600)  # cert_code -> pdf_bytes

import subprocess
import shutil

# Cap concurrent soffice processes — tune based on container CPU/RAM
# allocated in Dokploy (start with 2 for a 1-2 vCPU container)
LIBREOFFICE_SEMAPHORE = asyncio.Semaphore(2)

def export_pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
    """
    Converts PPTX to PDF using headless LibreOffice.
    Uses an isolated user profile per call to avoid lock-file
    collisions under concurrent requests.
    """
    os.makedirs(output_dir, exist_ok=True)

    # Isolated profile dir per invocation — prevents concurrent
    # soffice processes from fighting over the same lock file
    profile_dir = f"/tmp/lo_profile_{uuid.uuid4().hex}"

    cmd = [
        "soffice",
        "--headless",
        "--norestore",
        "--nofirststartwizard",
        "--nologo",
        "--nodefault",
        "--invisible",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        f"-env:UserInstallation=file://{profile_dir}",
        pptx_path,
    ]

    # Explicitly set HOME environment variable to /tmp to prevent LibreOffice
    # from failing when trying to write to non-writable Nixpacks/Dokploy directories.
    env = os.environ.copy()
    env["HOME"] = "/tmp"

    try:
        subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=60,  # hard timeout — do not let a stuck soffice process hang
            check=True,
            env=env,
        )
    except subprocess.TimeoutExpired:
        raise RuntimeError(
            f"LibreOffice conversion timed out for {pptx_path}"
        )
    except subprocess.CalledProcessError as e:
        raise RuntimeError(
            f"LibreOffice conversion failed with exit code {e.returncode}.\nStdout: {e.stdout}\nStderr: {e.stderr}"
        )
    finally:
        # Always clean up the temp profile directory, even on failure
        shutil.rmtree(profile_dir, ignore_errors=True)

    expected_pdf = os.path.join(
        output_dir,
        os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf",
    )

    if not os.path.exists(expected_pdf):
        raise RuntimeError(
            f"Conversion reported success but output PDF not found: "
            f"{expected_pdf}"
        )

    return expected_pdf

async def export_pptx_to_pdf_async(pptx_path: str, output_dir: str) -> str:
    async with LIBREOFFICE_SEMAPHORE:
        loop = asyncio.get_event_loop()
        return await loop.run_in_executor(
            None, export_pptx_to_pdf, pptx_path, output_dir
        )

TEMPLATE_EXPORT_LOCK = asyncio.Lock()

def get_alignment_str(paragraph):
    from pptx.enum.text import PP_ALIGN
    align = paragraph.alignment
    if align == PP_ALIGN.CENTER:
        return "center"
    elif align == PP_ALIGN.RIGHT:
        return "right"
    elif align == PP_ALIGN.JUSTIFY:
        return "justify"
    return "left"

def get_color_hex(run):
    try:
        color = run.font.color
        if color and color.type == 1:  # RGBColor
            rgb = color.rgb
            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
    except:
        pass
    return None

def is_body_text_shape(shape):
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    if not text:
        return False
    # Always include if it contains a placeholder
    if any(p in text for p in ["<<", ">>", "«", "»"]):
        return True
    # Bounding check for vertical column shapes (wide text blocks in the middle of A4)
    left_in = shape.left.inches
    top_in = shape.top.inches
    width_in = shape.width.inches
    
    is_wide_and_centered = (left_in < 1.0) and (width_in > 6.5)
    is_in_middle_band = (top_in > 2.2) and (top_in < 9.0)
    
    return is_wide_and_centered and is_in_middle_band

async def get_or_create_html_template(batch_id: str, cert_type: str, template_bytes: bytes) -> str:
    # Normalize batch_id to make it safe as a directory name
    safe_batch = re.sub(r'[\\/*?:"<>|\s]', "_", batch_id.strip()).lower()
    cache_dir = os.path.abspath(f"templates_cache/{safe_batch}/{cert_type}")
    
    layout_json_path = os.path.join(cache_dir, "layout.json")
    bg_png_path = os.path.join(cache_dir, "background.png")
    
    if os.path.exists(layout_json_path) and os.path.exists(bg_png_path):
        return cache_dir
        
    async with TEMPLATE_EXPORT_LOCK:
        # Double check inside the lock
        if os.path.exists(layout_json_path) and os.path.exists(bg_png_path):
            return cache_dir
            
        os.makedirs(cache_dir, exist_ok=True)
        print(f"[pdf] On-the-fly exporting template layout & background for {batch_id}/{cert_type}...")
        
        # 1. Save original pptx
        temp_pptx_path = os.path.join(cache_dir, "temp_original.pptx")
        with open(temp_pptx_path, "wb") as f:
            f.write(template_bytes)
            
        try:
            # 2. Parse shapes layout
            prs = Presentation(temp_pptx_path)
            slide = prs.slides[0]
            
            layout_data = {
                "template": cert_type,
                "width_in": prs.slide_width.inches,
                "height_in": prs.slide_height.inches,
                "shapes": []
            }
            
            shapes_to_clear = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if not is_body_text_shape(shape):
                    continue
                    
                left_in = shape.left.inches
                top_in = shape.top.inches
                width_in = shape.width.inches
                height_in = shape.height.inches
                
                font_name = "Calibri"
                font_size = 14
                font_color = "#000000"
                bold = False
                italic = False
                align = "left"
                
                if shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                    align = get_alignment_str(p)
                    if p.runs:
                        r = p.runs[0]
                        if r.font.name:
                            font_name = r.font.name
                        if r.font.size:
                            font_size = r.font.size.pt
                        c = get_color_hex(r)
                        if c:
                            font_color = c
                        bold = bool(r.font.bold)
                        italic = bool(r.font.italic)
                        
                paragraphs_cfg = []
                for p in shape.text_frame.paragraphs:
                    runs_cfg = []
                    for r in p.runs:
                        r_color = get_color_hex(r) or font_color
                        runs_cfg.append({
                            "text": r.text,
                            "font_name": r.font.name or font_name,
                            "font_size": r.font.size.pt if r.font.size else font_size,
                            "bold": bool(r.font.bold),
                            "italic": bool(r.font.italic),
                            "color": r_color
                        })
                    paragraphs_cfg.append({
                        "align": get_alignment_str(p),
                        "runs": runs_cfg
                    })
                    
                shape_cfg = {
                    "id": shape.shape_id,
                    "name": shape.name,
                    "left": left_in,
                    "top": top_in,
                    "width": width_in,
                    "height": height_in,
                    "font_name": font_name,
                    "font_size": font_size,
                    "color": font_color,
                    "bold": bold,
                    "italic": italic,
                    "align": align,
                    "original_text": shape.text_frame.text,
                    "paragraphs": paragraphs_cfg
                }
                layout_data["shapes"].append(shape_cfg)
                shapes_to_clear.append(shape)
                
            with open(layout_json_path, "w", encoding="utf-8") as f:
                json.dump(layout_data, f, indent=2)
                
            # 3. Clear text runs for background export
            for shape in shapes_to_clear:
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        r.text = ""
                        
            temp_cleared_pptx = os.path.join(cache_dir, "temp_cleared.pptx")
            prs.save(temp_cleared_pptx)
            
            # 4. Export background PNG (Windows COM or Linux LibreOffice)
            exported_bg = False
            try:
                import sys
                if sys.platform == "win32":
                    import win32com.client
                    print(f"  [Windows] Exporting background slide via PowerPoint COM...")
                    powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
                    presentation = powerpoint.Presentations.Open(os.path.abspath(temp_cleared_pptx), WithWindow=False)
                    presentation.Slides(1).Export(os.path.abspath(bg_png_path), "PNG")
                    presentation.Close()
                    powerpoint.Quit()
                    exported_bg = True
            except Exception as win_err:
                print(f"  [Windows] COM export failed/not available: {win_err}")
                
            if not exported_bg:
                # Fallback to LibreOffice + PyMuPDF
                pdf_path = await export_pptx_to_pdf_async(temp_cleared_pptx, cache_dir)
                doc = fitz.open(pdf_path)
                page = doc[0]
                matrix = fitz.Matrix(300 / 72, 300 / 72)
                pix = page.get_pixmap(matrix=matrix)
                pix.save(bg_png_path)
                doc.close()
                if os.path.exists(pdf_path):
                    os.remove(pdf_path)
                    
            # Cleanup temp cleared pptx
            if os.path.exists(temp_cleared_pptx):
                os.remove(temp_cleared_pptx)
                
            print(f"[pdf] On-the-fly export completed for {batch_id}/{cert_type}!")
        finally:
            if os.path.exists(temp_pptx_path):
                os.remove(temp_pptx_path)
                
        return cache_dir

def compile_weasyprint_pdf(html_content: str) -> bytes:
    from weasyprint import HTML
    return HTML(string=html_content).write_pdf()

async def generate_pdf_bytes_with_weasyprint_async(batch_id: str, cert_type: str, template_bytes: bytes, replacements: dict, qr_bytes: bytes) -> bytes:
    # 1. Get or create cache dir
    cache_dir = await get_or_create_html_template(batch_id, cert_type, template_bytes)
    
    # 2. Render layout HTML
    engine = LayoutEngine(cache_dir)
    html_content = engine.render_html(replacements, qr_bytes)
    
    # 3. Compile to PDF bytes asynchronously
    loop = asyncio.get_event_loop()
    pdf_bytes = await loop.run_in_executor(None, compile_weasyprint_pdf, html_content)
    return pdf_bytes

async def convert_pptx_to_pdf_bytes_async(pptx_bytes: bytes) -> bytes:
    temp_pptx = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    temp_pptx.write(pptx_bytes)
    temp_pptx.close()
    
    out_dir = os.path.dirname(os.path.abspath(temp_pptx.name))
    try:
        pdf_path = await export_pptx_to_pdf_async(os.path.abspath(temp_pptx.name), out_dir)
        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()
        try:
            os.remove(pdf_path)
        except Exception:
            pass
        return pdf_bytes
    finally:
        try:
            os.remove(temp_pptx.name)
        except Exception:
            pass



app = FastAPI(title="Certificate Generator API")

VERIFY_BASE_URL = os.getenv("CERT_VERIFY_BASE_URL", "https://coptercode-website.vercel.app/verify").strip()


def build_verify_url(cert_code: str) -> str:
    """
    Builds a verification URL for QR generation.
    Supports either:
    - direct placeholder format: https://site/verify/{code}
    - query format: https://site/verify?id=CODE
    """
    if not VERIFY_BASE_URL:
        return f"https://coptercode-website.vercel.app/verify?id={cert_code}"

    base = VERIFY_BASE_URL.rstrip('/')

    if "{code}" in base:
        return base.replace("{code}", cert_code)

    separator = "&" if "?" in base else "?"
    return f"{base}{separator}id={cert_code}"


BACKEND_BASE_URL = os.getenv("BACKEND_BASE_URL", "http://localhost:5000").strip()


async def send_email_notification(
    to_email: str,
    intern_name: str,
    batch_title: str,
    certificates: List[Dict[str, str]],
    intern_id: str
):
    import requests
    from email.mime.application import MIMEApplication

    smtp_server = os.getenv("SMTP_SERVER")
    smtp_port = os.getenv("SMTP_PORT", "587")
    smtp_username = os.getenv("SMTP_USERNAME")
    smtp_password = os.getenv("SMTP_PASSWORD")
    smtp_from_email = os.getenv("SMTP_FROM_EMAIL")
    smtp_from_name = os.getenv("SMTP_FROM_NAME", "CopterCode Team")

    intern_portal_base_url = os.getenv("INTERN_PORTAL_BASE_URL", "https://coptercode-website.vercel.app/intern-portal").strip().rstrip("/")
    email_logo_url = os.getenv("EMAIL_LOGO_URL", "https://coptercode-website.vercel.app/coptercode-logo.svg").strip()
    email_hero_image_url = os.getenv("EMAIL_HERO_IMAGE_URL", "https://coptercode-website.vercel.app/hero-3.jpg").strip()

    if not smtp_server or not smtp_username or not smtp_password or not smtp_from_email:
        print("WARNING: SMTP credentials not fully configured. Skipping email dispatch.")
        return False

    try:
        portal_link = f"{intern_portal_base_url}?id={intern_id}"

        # Use MIMEMultipart("mixed") since we are sending inline HTML + physical file attachments
        msg = MIMEMultipart("mixed")
        msg["Subject"] = "Congratulations on Successfully Completing Your Internship at CopterCode"
        msg["From"] = f"{smtp_from_name} <{smtp_from_email}>"
        msg["To"] = to_email

        # Create alternative part for HTML body
        msg_alternative = MIMEMultipart("alternative")
        msg.attach(msg_alternative)

        cert_items_html = ""
        for cert in certificates:
            cert_items_html += f"<li style='margin-bottom: 6px;'><strong>{cert['label']}</strong></li>"

        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>Congratulations from CopterCode</title>
        </head>
        <body style="font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Helvetica, Arial, sans-serif; background-color: #f8fafc; color: #1e293b; margin: 0; padding: 0; -webkit-font-smoothing: antialiased;">
            <table border="0" cellpadding="0" cellspacing="0" width="100%" style="background-color: #f8fafc; padding: 40px 0;">
                <tr>
                    <td align="center">
                        <table border="0" cellpadding="0" cellspacing="0" width="600" style="background-color: #ffffff; border-radius: 16px; overflow: hidden; box-shadow: 0 4px 6px -1px rgba(0,0,0,0.05), 0 2px 4px -1px rgba(0,0,0,0.035); border: 1px solid #e2e8f0;">
                            
                            <!-- Header / Logo side-by-side -->
                            <tr>
                                <td align="center" style="background-color: #ffffff; padding: 36px 24px 28px 24px;">
                                    <table border="0" cellpadding="0" cellspacing="0" align="center" style="margin: 0 auto;">
                                        <tr>
                                            <td style="vertical-align: middle;">
                                                <img src="{email_logo_url}" alt="CopterCode Logo" style="height: 48px; width: auto; display: block; border-radius: 12px;" />
                                            </td>
                                            <td style="vertical-align: middle; padding-left: 12px; font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; font-size: 26px; font-weight: 750; color: #0f172a; letter-spacing: -0.5px; line-height: 1;">
                                                CopterCode
                                            </td>
                                        </tr>
                                    </table>
                                </td>
                            </tr>
                            
                            <!-- Hero Image -->
                            <tr>
                                <td style="padding: 0 24px;">
                                    <img src="{email_hero_image_url}" alt="CopterCode Team" style="width: 100%; height: auto; display: block; border-radius: 12px; object-fit: cover;" />
                                </td>
                            </tr>
                            
                            <!-- Main Content -->
                            <tr>
                                <td style="padding: 36px 36px 24px 36px;">
                                    <h2 style="font-size: 16px; font-weight: 700; color: #0f172a; margin: 0 0 16px 0;">Greetings from CopterCode!</h2>
                                    
                                    <p style="font-size: 14px; line-height: 1.6; color: #475569; margin: 0 0 18px 0;">
                                        We are pleased to inform you that you have successfully completed the one-month internship program at <strong>CopterCode</strong>. Your dedication, hard work, and enthusiasm throughout the program have been truly commendable.
                                    </p>
                                    
                                    <p style="font-size: 14px; line-height: 1.6; color: #475569; margin: 0 0 18px 0;">
                                        As a token of your accomplishments, the following documents have been issued for your reference and future endeavors:
                                    </p>
                                    
                                    <!-- Dynamic list inside a premium left-bordered box -->
                                    <table border="0" cellpadding="0" cellspacing="0" width="100%" style="background-color: #f8fafc; border-top: 1px solid #e2e8f0; border-right: 1px solid #e2e8f0; border-bottom: 1px solid #e2e8f0; border-left: 4px solid #5844e9; border-radius: 8px; margin-bottom: 24px;">
                                        <tr>
                                            <td style="padding: 16px 20px 10px 20px;">
                                                <ul style="margin: 0; padding: 0 0 0 20px; font-size: 14px; line-height: 1.8; color: #334155;">
                                                    {cert_items_html}
                                                </ul>
                                            </td>
                                        </tr>
                                    </table>
                                    
                                    <p style="font-size: 14px; line-height: 1.6; color: #475569; margin: 24px 0 16px 0;">
                                        You can download, view, verify, and share your certificate(s) at any time by visiting your Intern Portal or by clicking the button below:
                                    </p>
                                    
                                    <!-- Download Certificates CTA Button -->
                                    <table border="0" cellpadding="0" cellspacing="0" width="100%" style="margin: 24px 0;">
                                        <tr>
                                            <td align="center">
                                                <table border="0" cellpadding="0" cellspacing="0" style="margin: 0 auto;">
                                                    <tr>
                                                        <td align="center" style="background-color: #5844e9; border-radius: 8px;">
                                                            <a href="{portal_link}" target="_blank" style="display: inline-block; padding: 14px 28px; font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif; font-size: 15px; font-weight: 700; color: #ffffff !important; text-decoration: none; border-radius: 8px;">Download Certificates</a>
                                                        </td>
                                                    </tr>
                                                </table>
                                            </td>
                                        </tr>
                                    </table>
                                    
                                    <p style="font-size: 14px; line-height: 1.6; color: #475569; margin: 24px 0 18px 0;">
                                        We extend our heartfelt congratulations and best wishes for your future career. We are confident that the skills and knowledge you have gained here will serve you well in all your professional pursuits.
                                    </p>
                                    
                                    <p style="font-size: 14px; line-height: 1.6; color: #475569; margin: 0 0 20px 0;">
                                        Please feel free to stay in touch with us for any guidance or opportunities. We look forward to seeing you achieve great success ahead.
                                    </p>
                                    
                                    <p style="margin: 24px 0 0 0; font-size: 14px; color: #475569; font-weight: 500; line-height: 1.5;">
                                        Best regards,<br>
                                        <span style="font-weight: 700; font-size: 15px; color: #0f172a; display: block; margin: 4px 0 2px 0;">HR Team</span>
                                        <span style="font-weight: 600; color: #334155;">CopterCode</span>
                                    </p>
                                </td>
                            </tr>
                            
                            <!-- Premium B&W Signature & Footer -->
                            <tr>
                                <td style="background-color: #0f172a; padding: 24px 36px; color: #f8fafc; font-size: 13px; line-height: 1.6; border-bottom-left-radius: 15px; border-bottom-right-radius: 15px;">
                                    <table border="0" cellpadding="0" cellspacing="0" width="100%">
                                        <tr>
                                            <!-- Left Signature details -->
                                            <td valign="top" align="left">
                                                <!-- Horizontal links row -->
                                                <table cellpadding="0" cellspacing="0" border="0" style="border-collapse: collapse;">
                                                    <tr>
                                                        <td style="padding-right: 16px; font-size: 13px; color: #cbd5e1;">
                                                            <a href="mailto:hr@coptercode.co.in" style="color: #a5b4fc !important; text-decoration: none; font-weight: 500;">hr@coptercode.co.in</a>
                                                        </td>
                                                        <td style="padding-right: 16px; font-size: 13px; color: #334155;">|</td>
                                                        <td style="padding-right: 16px; font-size: 13px; color: #cbd5e1;">
                                                            <a href="https://www.coptercode.co.in/" target="_blank" style="color: #a5b4fc !important; text-decoration: none; font-weight: 500;">https://www.coptercode.co.in/</a>
                                                        </td>
                                                        <td style="padding-right: 16px; font-size: 13px; color: #334155;">|</td>
                                                        <td style="padding-right: 16px; font-size: 13px; color: #cbd5e1;">
                                                            <a href="https://www.instagram.com/coptercode/" target="_blank" style="color: #a5b4fc !important; text-decoration: none; font-weight: 500;">Instagram</a>
                                                        </td>
                                                        <td style="padding-right: 16px; font-size: 13px; color: #334155;">|</td>
                                                        <td style="font-size: 13px; color: #cbd5e1;">
                                                            <a href="https://www.linkedin.com/company/coptercode/" target="_blank" style="color: #a5b4fc !important; text-decoration: none; font-weight: 500;">LinkedIn</a>
                                                        </td>
                                                    </tr>
                                                </table>
                                            </td>
                                        </tr>
                                        <tr>
                                            <td align="center" style="padding-top: 20px; border-top: 1px solid #334155; margin-top: 0; font-size: 10px; color: #64748b;">
                                                This is an automated message. Please do not reply directly to this mail.
                                            </td>
                                        </tr>
                                    </table>
                                </td>
                            </tr>
                        </table>
                    </td>
                </tr>
            </table>
        </body>
        </html>
        """

        msg_alternative.attach(MIMEText(html_content, "html"))

        # Generate and attach certificates dynamically in memory
        for cert in certificates:
            try:
                pdf_bytes = await get_pdf_bytes_for_certificate(cert["cert_code"])
                part = MIMEApplication(pdf_bytes, _subtype="pdf")
                safe_name = re.sub(r'[\\/*?:"<>|]', "_", str(intern_name)).strip()
                part.add_header(
                    "Content-Disposition",
                    "attachment",
                    filename=f"{safe_name}_{cert['label']}.pdf"
                )
                msg.attach(part)
            except Exception as attachment_err:
                print(f"Failed to generate and attach PDF {cert['cert_code']} to email: {attachment_err}")

        server = smtplib.SMTP(smtp_server, int(smtp_port))
        if int(smtp_port) == 587:
            server.starttls()
        server.login(smtp_username, smtp_password)
        server.sendmail(smtp_from_email, to_email, msg.as_string())
        server.quit()
        print(f"Email sent successfully to {to_email}")
        return True
    except Exception as email_err:
        print(f"Error sending email to {to_email}: {email_err}")
        return False

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000",
        "http://localhost:5000",
        "https://coptercode-certificate.vercel.app",
        "https://coptercode-website.vercel.app",
        "https://coptercode.co.in",
        "https://cert-gen.coptercode.co.in"
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
def read_root():
    return {"status": "success", "message": "Certificate Generator API is running"}

# Supabase credentials verification
SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_SERVICE_ROLE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY")

if not SUPABASE_URL or not SUPABASE_SERVICE_ROLE_KEY:
    print("WARNING: SUPABASE_URL or SUPABASE_SERVICE_ROLE_KEY environment variables are missing.")

# Initialize Supabase client
supabase: Optional[Client] = None
if SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY:
    try:
        supabase = create_client(SUPABASE_URL, SUPABASE_SERVICE_ROLE_KEY)
    except Exception as e:
        print(f"Error initializing Supabase client: {e}")

# ───────────────────────────────────────────────────────────────────
# Document conversion
# ───────────────────────────────────────────────────────────────────

def convert_html_to_pdf_bytes(html_content: str) -> bytes:
    import io
    from xhtml2pdf import pisa
    
    pdf_buffer = io.BytesIO()
    pisa_status = pisa.CreatePDF(html_content, dest=pdf_buffer)
    if pisa_status.err:
        raise ValueError("Failed to convert HTML to PDF via xhtml2pdf")
    return pdf_buffer.getvalue()


def generate_certificate_from_html_bytes(html_bytes: bytes, replacements: dict, qr_bytes: bytes) -> tuple:
    import base64
    
    # Decode the HTML template
    html_text = html_bytes.decode("utf-8", errors="ignore")
    
    # Generate base64 for QR Code
    qr_base64 = base64.b64encode(qr_bytes).decode("utf-8")
    qr_img_tag = f"data:image/png;base64,{qr_base64}"
    
    # Standardize replacements to handle both <<KEY>> and {{KEY}} formats
    for key, val in list(replacements.items()):
        # Replace <<KEY>>
        html_text = html_text.replace(key, val)
        # Also replace curly brace format e.g. {{KEY}}
        curly_key = key.replace("<<", "{{").replace(">>", "}}")
        html_text = html_text.replace(curly_key, val)
        # Also replace French quotes format e.g. «KEY»
        guill_key = key.replace("<<", "«").replace(">>", "»")
        html_text = html_text.replace(guill_key, val)
        
    # Replace QR placeholder
    for qr_key in ("<<QR>>", "<<QR_CODE>>", "<<QRCODE>>", "{{QR}}", "{{QR_CODE}}", "{{QRCODE}}", "«QR»", "«QR_CODE»", "«QRCODE»"):
        html_text = html_text.replace(qr_key, qr_img_tag)
        
    # Detect certificate title based on keywords
    html_lower = html_text.lower()
    if "recommendation" in html_lower or "recomandation" in html_lower:
        cert_title = "Letter Of Recomandation"
    elif "experience certificate" in html_lower or "experience letter" in html_lower:
        cert_title = "Experience Letter"
    elif "internship" in html_lower:
        cert_title = "Internship Certificate"
    else:
        cert_title = "Certificate"
        
    pdf_bytes = convert_html_to_pdf_bytes(html_text)
    return pdf_bytes, cert_title


def convert_pptx_to_pdf_bytes(pptx_bytes: bytes) -> bytes:
    temp_pptx = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    temp_pptx.write(pptx_bytes)
    temp_pptx.close()
    
    out_dir = os.path.dirname(os.path.abspath(temp_pptx.name))
    try:
        pdf_path = export_pptx_to_pdf(os.path.abspath(temp_pptx.name), out_dir)
        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()
        try:
            os.remove(pdf_path)
        except Exception:
            pass
        return pdf_bytes
    finally:
        try:
            os.remove(temp_pptx.name)
        except Exception:
            pass


def defragment_paragraph(paragraph, placeholders):
    """
    Defragments python-pptx runs that have split placeholders
    (e.g., <<INTERNSHIP & LIVE PROJECT AREA>> might be split into multiple runs).
    Converts split runs into a single run so they can be replaced accurately.
    """
    for key in placeholders:
        keys_to_check = [key, key.replace("<<", "«").replace(">>", "»")]
        for k_check in keys_to_check:
            search_start = 0
            while True:
                p_text = "".join(run.text for run in paragraph.runs)
                start_idx = p_text.find(k_check, search_start)
                if start_idx == -1:
                    break

                # Map each character position to its run index
                run_map = []
                for run_idx, run in enumerate(paragraph.runs):
                    run_map.extend([run_idx] * len(run.text))

                end_idx = start_idx + len(k_check) - 1
                if end_idx >= len(run_map):
                    break

                # Get run indices that cover the placeholder
                run_indices = run_map[start_idx : end_idx + 1]
                first_run_idx = run_indices[0]
                last_run_idx = run_indices[-1]

                if first_run_idx == last_run_idx:
                    search_start = start_idx + len(k_check)
                    continue

                # Merge split runs
                combined_text = "".join(paragraph.runs[r_i].text for r_i in range(first_run_idx, last_run_idx + 1))
                paragraph.runs[first_run_idx].text = combined_text
                for r_i in range(first_run_idx + 1, last_run_idx + 1):
                    paragraph.runs[r_i].text = ""

                search_start = start_idx + len(k_check)

def text_frame_has_placeholders(text_frame, replacements) -> bool:
    """
    Returns True if the text frame contains any of the keys in replacements,
    or general QR placeholders.
    """
    text = text_frame.text
    for key in replacements:
        if key in text:
            return True
        guill_key = key.replace("<<", "«").replace(">>", "»")
        if guill_key in text:
            return True
    
    for qr_key in ("<<QR>>", "<<QR_CODE>>", "<<QRCODE>>", "«QR»", "«QR_CODE»", "«QRCODE»"):
        if qr_key in text:
            return True
            
    return False

def detect_certificate_title(slide) -> str:
    """
    Detects the certificate type based on the text contents of the slide shapes.
    Returns: 'Internship Certificate', 'Letter Of Recomandation', 'Experience Letter', or 'Certificate'.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if "recommendation" in text or "recomandation" in text:
                return "Letter Of Recomandation"
            elif "experience certificate" in text or "experience letter" in text:
                return "Experience Letter"
            elif "internship" in text:
                return "Internship Certificate"
    
    # Fallback checks
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if "certificate" in text:
                return "Certificate"
                
    return "Certificate"

def generate_certificate_from_pptx_bytes(pptx_bytes: bytes, replacements: dict, qr_bytes: bytes) -> tuple:
    from pptx import Presentation
    from pptx.util import Inches
    import tempfile
    
    # Save the input pptx bytes to a temp file
    temp_pptx = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    temp_pptx.write(pptx_bytes)
    temp_pptx.close()
    
    try:
        prs = Presentation(temp_pptx.name)
        slide = prs.slides[0]
        
        # Fields that should WRAP to next line instead of shrinking font
        WRAP_FIELDS = {"<<INTERNSHIP & LIVE PROJECT AREA>>", "<<PROJECT>>", "<<DOMAIN>>", "<<ROLE>>"}

        # 1. Process shapes and replace text
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                text_frame = shape.text_frame
                combined_text = text_frame.text.strip()

                # Fix for "COPTERCODE" wrapping "E" to the next line in the header textbox
                if ("COPTERCODE" in combined_text or "COPTERCOD" in combined_text) and len(combined_text) < 15:
                    text_frame.word_wrap = False
                    shape.width = shape.width + Inches(0.5)

                # Skip text boxes that don't contain any placeholders at all to preserve formatting/alignment
                if not text_frame_has_placeholders(text_frame, replacements):
                    continue

                # Enable word wrap so long values flow to next line naturally
                text_frame.word_wrap = True

                # Defragment paragraph runs to resolve split placeholders
                placeholders_to_defrag = list(replacements.keys()) + ["<<QR>>", "<<QR_CODE>>", "<<QRCODE>>"]
                for paragraph in text_frame.paragraphs:
                    defragment_paragraph(paragraph, placeholders_to_defrag)

                combined_text = text_frame.text.strip()

                # Check for QR placeholder
                if "<<QR>>" in combined_text or "«QR»" in combined_text:
                    # Make the QR code larger (1.15 in x 1.15 in)
                    qr_w = Inches(1.15)
                    qr_h = Inches(1.15)
                    
                    # If the placeholder textbox is abnormally wide, align the QR image to its left edge.
                    # Otherwise, center it.
                    if shape.width > Inches(2.5):
                        qr_left = shape.left
                    else:
                        qr_left = int(shape.left + (shape.width - qr_w) // 2)
                        
                    qr_top = int(shape.top + (shape.height - qr_h) // 2)

                    # Remove placeholder shape
                    sp = shape._element
                    sp.getparent().remove(sp)

                    # Write QR bytes to a temp file to insert as picture
                    temp_qr = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
                    temp_qr.write(qr_bytes)
                    temp_qr.close()

                    try:
                        slide.shapes.add_picture(temp_qr.name, qr_left, qr_top, width=qr_w, height=qr_h)
                    finally:
                        try:
                            os.remove(temp_qr.name)
                        except Exception:
                            pass
                    continue

                # Standard replacements on text runs
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, val in replacements.items():
                            has_match = False
                            actual_key = None

                            if key in run.text:
                                has_match = True
                                actual_key = key
                            else:
                                guill_key = key.replace("<<", "«").replace(">>", "»")
                                if guill_key in run.text:
                                    has_match = True
                                    actual_key = guill_key

                            if has_match:
                                # Standardize to Arial to prevent system font substitutions
                                run.font.name = "Arial"
                                
                                # For body/area fields: keep font size and let text wrap
                                # For short label fields (name, institution): allow mild shrinking
                                is_wrap_field = key in WRAP_FIELDS
                                if not is_wrap_field and run.font.size and val:
                                    length = len(val)
                                    limit = 20 if "NAME" in key else 30
                                    if length > limit:
                                        scale = limit / length
                                        scale = max(0.75, min(1.0, scale))
                                        run.font.size = int(run.font.size * scale)

                                run.text = run.text.replace(actual_key, val)

                                # Clean up visual hacks (multiple consecutive spaces) and fix missing comma spaces
                                if run.text:
                                    run.text = re.sub(r',([a-zA-Z])', r', \1', run.text)
                                    run.text = re.sub(r'\s{2,}', ' ', run.text)
                                
        # Save presentation back to same temp file
        prs.save(temp_pptx.name)
        
        # Detect certificate title from slide text
        cert_title = detect_certificate_title(slide)
        
        # Read the modified PPTX file bytes
        with open(temp_pptx.name, "rb") as f:
            modified_pptx_bytes = f.read()
            
        # Convert to PDF
        pdf_bytes = convert_pptx_to_pdf_bytes(modified_pptx_bytes)
        return pdf_bytes, cert_title
        
    finally:
        try:
            os.remove(temp_pptx.name)
        except Exception:
            pass


def render_certificate_pdf(template_bytes: bytes, replacements: dict, qr_bytes: bytes, ext: str = ".pptx") -> bytes:
    if ext.lower() == ".pptx":
        pdf_bytes, _ = generate_certificate_from_pptx_bytes(template_bytes, replacements, qr_bytes)
    else:
        pdf_bytes, _ = generate_certificate_from_html_bytes(template_bytes, replacements, qr_bytes)
    return pdf_bytes


def upload_pdf_to_storage(cert_code: str, pdf_bytes: bytes) -> str:
    """
    Uploads the generated PDF to Supabase Storage bucket 'certificates'
    and returns the public URL.
    """
    if not supabase:
        raise ValueError("Supabase client is not initialized.")
    
    bucket_name = "certificates"
    path = f"{cert_code}.pdf"
    
    try:
        supabase.storage.from_(bucket_name).upload(
            path=path,
            file=pdf_bytes,
            file_options={"content-type": "application/pdf", "upsert": "true"}
        )
    except Exception as e:
        print(f"[storage] Upload failed: {e}. Trying to create bucket '{bucket_name}'...")
        try:
            supabase.storage.create_bucket(bucket_name, options={"public": True})
            supabase.storage.from_(bucket_name).upload(
                path=path,
                file=pdf_bytes,
                file_options={"content-type": "application/pdf", "upsert": "true"}
            )
        except Exception as retry_err:
            print(f"[storage] Retry upload failed: {retry_err}")
            raise retry_err

    public_url = supabase.storage.from_(bucket_name).get_public_url(path)
    return public_url


def generate_and_store_pdf(cert_code: str, cert_type: str, intern_data: dict, batch_id: str) -> str:
    """
    Generates the certificate PDF and uploads it to Supabase Storage.
    Returns the public URL of the uploaded PDF.
    """
    if not supabase:
        raise ValueError("Supabase client is not configured.")

    replacements = {
        "<<NAME>>": intern_data.get("name") or "",
        "<<INSTITUTION>>": intern_data.get("college") or "",
        "<<COLLEGE>>": intern_data.get("college") or "",
        "<<YEAR>>": intern_data.get("year") or "",
        "<<DEPARTMENT>>": intern_data.get("department") or "",
        "<<DOMAIN>>": intern_data.get("role") or "",
        "<<ROLE>>": intern_data.get("role") or "",
        "<<PROJECT>>": intern_data.get("project") or "",
        "<<INTERNSHIP & LIVE PROJECT AREA>>": intern_data.get("project") or "",
        "<<BATCH>>": intern_data.get("month") or "",
        "<<BATCH >>": intern_data.get("month") or "",
        "<<DATE>>": str(intern_data.get("date")) if intern_data.get("date") else "",
        "<<DT>>": str(intern_data.get("date")) if intern_data.get("date") else ""
    }

    qr_url = build_verify_url(cert_code)
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=12,
        border=2
    )
    qr.add_data(qr_url)
    qr.make(fit=True)
    from qrcode.image.pil import PilImage
    qr_img = qr.make_image(image_factory=PilImage, fill_color="black", back_color="white")
    qr_io = io.BytesIO()
    try:
        qr_img.save(qr_io, format="PNG")
    except TypeError:
        qr_io.seek(0)
        qr_io.truncate(0)
        qr_img.save(qr_io)
    qr_bytes = qr_io.getvalue()

    # Resolve the .pptx template directly — skip the overlay pipeline entirely.
    template_bytes = None
    ext = ".pptx"

    # Priority 1: named file e.g. templates/BATCH_TEST1/internship.pptx
    if batch_id and not batch_id.startswith("slot-"):
        import re as _re
        batch_id_stripped = batch_id.strip()
        batch_id_normalized = _re.sub(r'\s+', '_', batch_id_stripped.lower())
        cert_type_file = f"{cert_type}.pptx"

        for folder in [batch_id_stripped, batch_id_normalized]:
            if not folder:
                continue
            try:
                candidate = f"templates/{folder}/{cert_type_file}"
                template_bytes = supabase.storage.from_("templates").download(candidate)
                print(f"[pdf] Using named PPTX: {candidate}")
                break
            except Exception:
                pass

    # Priority 2: slot-based path e.g. templates/slot-1.pptx
    if not template_bytes:
        slot_path = f"templates/{batch_id}.pptx"
        try:
            template_bytes = supabase.storage.from_("templates").download(slot_path)
            print(f"[pdf] Using slot PPTX: {slot_path}")
        except Exception:
            pass

    # Priority 3: batches table lookup
    if not template_bytes:
        try:
            batch_res = supabase.table("batches").select("*").eq("id", batch_id).execute()
            if batch_res.data:
                batch = batch_res.data[0]
                if cert_type == "lor":
                    template_path = batch.get("lor_template_path")
                elif cert_type == "experience":
                    template_path = batch.get("experience_template_path")
                else:
                    template_path = batch.get("internship_template_path")
                if template_path:
                    ext = os.path.splitext(template_path)[1].lower() or ".pptx"
                    template_bytes = supabase.storage.from_("templates").download(template_path)
                    print(f"[pdf] Using batches-table path: {template_path}")
        except Exception:
            pass

    if not template_bytes:
        raise ValueError(f"Template PPTX not found for batch '{batch_id}' / cert_type '{cert_type}'. "
                         f"Upload templates/{{batch_id}}/{{cert_type}}.pptx to storage.")

    if ext == ".pptx":
        modified_pptx_bytes, _ = generate_certificate_from_pptx_bytes(template_bytes, replacements, qr_bytes)
        pdf_bytes = convert_pptx_to_pdf_bytes(modified_pptx_bytes)
    else:
        pdf_bytes, _ = generate_certificate_from_html_bytes(template_bytes, replacements, qr_bytes)

    public_url = upload_pdf_to_storage(cert_code, pdf_bytes)
    return public_url



def has_placeholders(text: str) -> bool:
    if not text:
        return False
    # Matches <<...>>, {{...}}, or «...»
    pattern = r'(<<.*?>>|\{\{.*?\}\}|«.*?»)'
    return bool(re.search(pattern, text))


def get_pptx_alignment_name(alignment) -> str:
    # PP_ALIGN enum values:
    # LEFT = 1, CENTER = 2, RIGHT = 3, JUSTIFY = 4
    if alignment == 1:
        return "left"
    elif alignment == 2:
        return "center"
    elif alignment == 3:
        return "right"
    elif alignment == 4:
        return "justify"
    return "center"


def get_font_properties(paragraph, run=None):
    font = run.font if (run and run.font) else paragraph.font
    
    # Font name
    font_name = font.name
    if not font_name and paragraph.font:
        font_name = paragraph.font.name
    if not font_name:
        font_name = "Arial"
        
    # Font size
    font_size = None
    if font.size:
        font_size = font.size.pt
    elif paragraph.font and paragraph.font.size:
        font_size = paragraph.font.size.pt
    if not font_size:
        font_size = 12.0
        
    # Bold
    bold = font.bold
    if bold is None and paragraph.font:
        bold = paragraph.font.bold
    bold = bool(bold)
    
    # Color
    color_hex = "#000000"
    try:
        if font.color and font.color.type == 1:
            color_hex = f"#{font.color.rgb[0]:02x}{font.color.rgb[1]:02x}{font.color.rgb[2]:02x}"
        elif paragraph.font and paragraph.font.color and paragraph.font.color.type == 1:
            color_hex = f"#{paragraph.font.color.rgb[0]:02x}{paragraph.font.color.rgb[1]:02x}{paragraph.font.color.rgb[2]:02x}"
    except Exception:
        pass
        
    return font_name, font_size, bold, color_hex


def get_paragraph_html(paragraph) -> str:
    if not paragraph.runs:
        return paragraph.text.replace("&", "&amp;")
    html_parts = []
    for run in paragraph.runs:
        text = run.text
        if not text:
            continue
        escaped_text = text.replace("&", "&amp;")
        if run.font.bold:
            escaped_text = f"<b>{escaped_text}</b>"
        if run.font.italic:
            escaped_text = f"<i>{escaped_text}</i>"
        if run.font.color and run.font.color.type == 1:
            c = run.font.color.rgb
            hex_color = f"#{c[0]:02x}{c[1]:02x}{c[2]:02x}"
            escaped_text = f'<font color="{hex_color}">{escaped_text}</font>'
        if run.font.name:
            escaped_text = f'<font name="{run.font.name}">{escaped_text}</font>'
        if run.font.size:
            escaped_text = f'<font size="{run.font.size.pt}">{escaped_text}</font>'
        html_parts.append(escaped_text)
    return "".join(html_parts)


def precompile_pptx_template(pptx_bytes: bytes) -> tuple:
    """
    Parses PPTX file, extracts all placeholder layouts, builds a JSON coordinate map,
    clears placeholder text in PPTX to blank them out, and returns:
    (blanked_pptx_bytes, coordinates_json_bytes)
    """
    import json
    from pptx import Presentation
    import tempfile
    
    temp_in = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    temp_in.write(pptx_bytes)
    temp_in.close()
    
    try:
        prs = Presentation(temp_in.name)
        slide = prs.slides[0]
        
        slide_width = prs.slide_width / 12700
        slide_height = prs.slide_height / 12700
        
        placeholders_metadata = []
        default_keys = [
            "<<NAME>>", "<<INSTITUTION>>", "<<COLLEGE>>", "<<YEAR>>", "<<DEPARTMENT>>",
            "<<DOMAIN>>", "<<ROLE>>", "<<PROJECT>>", "<<INTERNSHIP & LIVE PROJECT AREA>>",
            "<<BATCH>>", "<<BATCH >>", "<<DATE>>", "<<DT>>", "<<QR>>", "<<QR_CODE>>", "<<QRCODE>>"
        ]
        
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                text_frame = shape.text_frame
                combined_text = text_frame.text.strip()
                
                if not has_placeholders(combined_text):
                    continue
                
                # Defragment paragraph runs to resolve split placeholders
                for paragraph in text_frame.paragraphs:
                    defragment_paragraph(paragraph, default_keys)
                
                combined_text = text_frame.text.strip()
                
                left_pt = shape.left / 12700
                top_pt = shape.top / 12700
                width_pt = shape.width / 12700
                height_pt = shape.height / 12700
                
                is_qr = False
                for qr_key in ("<<QR>>", "<<QR_CODE>>", "<<QRCODE>>", "«QR»", "«QR_CODE»", "«QRCODE»"):
                    if qr_key in combined_text:
                        is_qr = True
                        break
                
                if is_qr:
                    placeholders_metadata.append({
                        "type": "qr",
                        "box": {
                            "left": left_pt,
                            "top": top_pt,
                            "width": width_pt,
                            "height": height_pt
                        }
                    })
                else:
                    paragraphs_meta = []
                    for paragraph in text_frame.paragraphs:
                        align = get_pptx_alignment_name(paragraph.alignment)
                        text_template = get_paragraph_html(paragraph)
                        font_name, font_size, font_bold, font_color = get_font_properties(paragraph, paragraph.runs[0] if paragraph.runs else None)
                        
                        paragraphs_meta.append({
                            "align": align,
                            "text_template": text_template,
                            "font_name": font_name,
                            "font_size": font_size,
                            "font_bold": font_bold,
                            "font_color": font_color
                        })
                        
                    from pptx.enum.text import MSO_ANCHOR
                    anchor = text_frame.vertical_anchor
                    v_align = "top"
                    if anchor == MSO_ANCHOR.MIDDLE:
                        v_align = "middle"
                    elif anchor == MSO_ANCHOR.BOTTOM:
                        v_align = "bottom"

                    placeholders_metadata.append({
                        "type": "text",
                        "box": {
                            "left": left_pt,
                            "top": top_pt,
                            "width": width_pt,
                            "height": height_pt
                        },
                        "vertical_anchor": v_align,
                        "paragraphs": paragraphs_meta
                    })
                
                # Blank out text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ""
                        
        temp_out = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        temp_out.close()
        prs.save(temp_out.name)
        
        with open(temp_out.name, "rb") as f:
            blanked_pptx_bytes = f.read()
            
        try:
            os.remove(temp_out.name)
        except Exception:
            pass
            
        coord_map = {
            "slide_width": slide_width,
            "slide_height": slide_height,
            "placeholders": placeholders_metadata
        }
        coordinates_json_bytes = json.dumps(coord_map, indent=2).encode("utf-8")
        
        return blanked_pptx_bytes, coordinates_json_bytes
        
    finally:
        try:
            os.remove(temp_in.name)
        except Exception:
            pass


def replace_tokens(template_str: str, replacements: dict) -> str:
    import html
    replaced = template_str
    for key, val in replacements.items():
        # Escape the replacement value so it is safe for ReportLab's HTML parser
        escaped_val = html.escape(str(val)) if val is not None else ""
        
        # Replace the raw key
        replaced = replaced.replace(key, escaped_val)
        
        # Replace the key with only ampersand escaped (matching get_paragraph_html)
        amp_key = key.replace("&", "&amp;")
        replaced = replaced.replace(amp_key, escaped_val)
        
        # Replace the fully HTML-escaped key
        escaped_key = html.escape(key)
        replaced = replaced.replace(escaped_key, escaped_val)
        
        # Also replace curly brace and French quote formatting
        for k in (key, amp_key, escaped_key):
            curly_key = k.replace("<<", "{{").replace(">>", "}}")
            replaced = replaced.replace(curly_key, escaped_val)
            guill_key = k.replace("<<", "«").replace(">>", "»")
            replaced = replaced.replace(guill_key, escaped_val)
            
    return replaced


def resolve_reportlab_font(font_name: str, font_bold: bool = False) -> str:
    font_name_lower = font_name.lower()
    if "canva sans" in font_name_lower or "canvasans" in font_name_lower:
        return "CanvaSans-Bold" if font_bold else "CanvaSans-Regular"
    elif "calibri" in font_name_lower:
        return "Calibri-Bold" if font_bold else "Calibri"
    elif "codec" in font_name_lower or "codecpro" in font_name_lower:
        return "CodecPro-Bold" if font_bold else "CodecPro-Regular"
    elif "arial" in font_name_lower or "helvetica" in font_name_lower or "sans" in font_name_lower:
        return "Helvetica-Bold" if font_bold else "Helvetica"
    else:
        if font_name in pdfmetrics.getRegisteredFontNames():
            return font_name
        elif font_bold and f"{font_name}-Bold" in pdfmetrics.getRegisteredFontNames():
            return f"{font_name}-Bold"
        else:
            return "Helvetica-Bold" if font_bold else "Helvetica"


def generate_overlay_pdf_bytes(
    background_pdf_bytes: bytes,
    coordinates_json_bytes: bytes,
    replacements: dict,
    qr_bytes: bytes
) -> bytes:
    import json
    import io
    from reportlab.pdfgen import canvas
    
    coords = json.loads(coordinates_json_bytes.decode("utf-8"))
    slide_width = coords.get("slide_width", 720.0)
    slide_height = coords.get("slide_height", 540.0)
    placeholders = coords.get("placeholders", [])
    
    # Extract matching font family and color from template text placeholders to match formatting
    text_font_name = "Arial"
    text_font_color = "#000000"
    for p_item in placeholders:
        if p_item.get("type") == "text":
            paragraphs = p_item.get("paragraphs", [])
            if paragraphs:
                text_font_name = paragraphs[0].get("font_name", "Arial")
                text_font_color = paragraphs[0].get("font_color", "#000000")
                break
    resolved_scan_font = resolve_reportlab_font(text_font_name, font_bold=False)
    
    overlay_io = io.BytesIO()
    can = canvas.Canvas(overlay_io, pagesize=(slide_width, slide_height))
    
    qr_reader = ImageReader(io.BytesIO(qr_bytes)) if qr_bytes else None
    
    for p in placeholders:
        p_type = p.get("type")
        box = p.get("box", {})
        left_pt = box.get("left", 0.0)
        top_pt = box.get("top", 0.0)
        width_pt = box.get("width", 0.0)
        height_pt = box.get("height", 0.0)
        
        bottom_pt = slide_height - (top_pt + height_pt)
        
        if p_type == "qr" and qr_reader:
            # Replicate PPTX QR sizing and alignment override logic to prevent stretched/distorted QR codes.
            # Standard size is 1.15 in = 82.8 pt
            qr_size = 82.8
            
            # Horizontally align:
            # If the placeholder textbox is abnormally wide (e.g. > 180 pt / 2.5 inches), align the QR image to its left edge.
            # Otherwise, center it horizontally within the textbox.
            if width_pt > 180.0:
                qr_left = left_pt
            else:
                qr_left = left_pt + (width_pt - qr_size) / 2
                
            # Vertically align:
            # Center the QR image vertically on the textbox line.
            qr_bottom = (slide_height - top_pt - height_pt / 2) - qr_size / 2
            
            can.drawImage(qr_reader, qr_left, qr_bottom, qr_size, qr_size)
            
            # Draw "Scan to Verify" text centered at the bottom of the QR code
            can.saveState()
            can.setFont(resolved_scan_font, 7.5) # Elegant matching font size
            from reportlab.lib.colors import HexColor
            can.setFillColor(HexColor(text_font_color))
            # Center of the QR code horizontally:
            scan_x = qr_left + qr_size / 2
            # Vertically below the bottom of the QR code:
            scan_y = qr_bottom - 11.0
            can.drawCentredString(scan_x, scan_y, "Scan to Verify")
            can.restoreState()
        elif p_type == "text":
            paragraphs = p.get("paragraphs", [])
            flowables = []
            for p_meta in paragraphs:
                text_template = p_meta.get("text_template", "")
                replaced_text = replace_tokens(text_template, replacements)
                
                # Clean visual spaces and comma formatting
                replaced_text = re.sub(r',([a-zA-Z])', r', \1', replaced_text)
                replaced_text = re.sub(r'\s{2,}', ' ', replaced_text)
                
                # Resolve nested <font face="..."> tags in HTML and map to correct attribute
                def replace_font_tag(match):
                    font_attr = match.group(1)
                    is_bold = "bold" in font_attr.lower()
                    resolved = resolve_reportlab_font(font_attr, is_bold)
                    return f'<font face="{resolved}">'
                
                replaced_text = re.sub(r'<font\s+(?:name|face)="([^"]+)">', replace_font_tag, replaced_text)
                
                align_map = {"left": TA_LEFT, "center": TA_CENTER, "right": TA_RIGHT, "justify": TA_JUSTIFY}
                alignment = align_map.get(p_meta.get("align"), TA_CENTER)
                
                font_name = p_meta.get("font_name", "Arial")
                font_size = p_meta.get("font_size", 12.0)
                font_bold = p_meta.get("font_bold", False)
                font_color = p_meta.get("font_color", "#000000")
                
                font_resolved = resolve_reportlab_font(font_name, font_bold)
                
                style = ParagraphStyle(
                    name=f"Style_{uuid.uuid4().hex[:8]}",
                    fontName=font_resolved,
                    fontSize=font_size,
                    leading=font_size * 1.15, # Tight leading to match PPTX line height
                    textColor=font_color,
                    alignment=alignment,
                    spaceBefore=0,
                    spaceAfter=0
                )
                
                flowables.append(Paragraph(replaced_text, style))
                
            # Perform vertical alignment adjustment based on PowerPoint anchor alignment
            v_align = p.get("vertical_anchor", "middle" if len(flowables) == 1 else "top")
            total_occupied_height = 0
            for f in flowables:
                _, h = f.wrap(width_pt, height_pt)
                total_occupied_height += h
                
            if total_occupied_height < height_pt:
                if v_align == "middle":
                    bottom_pt = bottom_pt + (height_pt - total_occupied_height) / 2
                    height_pt = total_occupied_height
                elif v_align == "bottom":
                    height_pt = total_occupied_height
                    
            frame = Frame(left_pt, bottom_pt, width_pt, height_pt,
                          leftPadding=0, rightPadding=0, topPadding=0, bottomPadding=0)
            
            kif = KeepInFrame(width_pt, height_pt, flowables, mode='shrink')
            frame.addFromList([kif], can)
            
    can.save()
    overlay_io.seek(0)
    
    bg_reader = PdfReader(io.BytesIO(background_pdf_bytes))
    overlay_reader = PdfReader(overlay_io)
    
    writer = PdfWriter()
    bg_page = bg_reader.pages[0]
    overlay_page = overlay_reader.pages[0]
    
    bg_page.merge_page(overlay_page)
    writer.add_page(bg_page)
    
    for i in range(1, len(bg_reader.pages)):
        writer.add_page(bg_reader.pages[i])
        
    out_pdf_io = io.BytesIO()
    writer.write(out_pdf_io)
    return out_pdf_io.getvalue()


@app.post("/api/templates/upload")
async def upload_template(
    slot: int = Form(...),
    template_file: UploadFile = File(...)
):
    if slot not in (1, 2, 3):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Slot must be 1, 2, or 3."
        )
        
    filename = template_file.filename.lower()
    ext = os.path.splitext(filename)[1]
    if ext != ".pptx":
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Template file must be a .pptx file."
        )
        
    try:
        pptx_bytes = await template_file.read()
        
        # Save to temp file to read with python-pptx
        temp_file = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        temp_file.write(pptx_bytes)
        temp_file.close()
        
        placeholders = set()
        try:
            from pptx import Presentation
            prs = Presentation(temp_file.name)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        found = re.findall(r'(<<.*?>>|\{\{.*?\}\}|«.*?»)', text)
                        for f in found:
                            placeholders.add(f)
        finally:
            try:
                os.remove(temp_file.name)
            except Exception:
                pass
                
        # Upload template to storage at templates/slot-{slot}.pptx, overwriting
        path = f"templates/slot-{slot}.pptx"
        content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        
        if not supabase:
            raise ValueError("Supabase client not configured.")
            
        supabase.storage.from_("templates").upload(
            path=path,
            file=pptx_bytes,
            file_options={"content-type": content_type, "upsert": "true"}
        )
        
        # Upsert template row in batches table
        batch_record = {
            "id": f"slot-{slot}",
            "month": f"Slot {slot} Template",
            "issue_date": datetime.utcnow().date().isoformat(),
            "internship_template_path": path
        }
        supabase.table("batches").upsert(batch_record).execute()
        
        # Cache the template bytes in memory
        template_cache[f"slot-{slot}"] = pptx_bytes
        
        return {
            "status": "success",
            "slot": slot,
            "placeholders": sorted(list(placeholders))
        }
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to upload and parse template: {str(e)}"
        )


@app.post("/api/batch/create")
async def create_batch(
    batch_id: str = Form(...),
    month: str = Form(...),
    issue_date: str = Form(...),
    lor_template: Optional[UploadFile] = File(None),
    experience_template: Optional[UploadFile] = File(None),
    internship_template: Optional[UploadFile] = File(None)
):
    batch_id = batch_id.strip()
    if not supabase:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Supabase client is not configured."
        )

    for f in (lor_template, experience_template, internship_template):
        if f:
            ext = os.path.splitext(f.filename)[1].lower()
            if ext not in (".html", ".pptx"):
                raise HTTPException(
                    status_code=status.HTTP_400_BAD_REQUEST,
                    detail=f"Template {f.filename} must be an HTML or PPTX file."
                )

    try:
        lor_path = None
        exp_path = None
        int_path = None

        if lor_template:
            lor_bytes = await lor_template.read()
            ext = os.path.splitext(lor_template.filename)[1].lower()
            lor_path = f"templates/{batch_id}/lor{ext}"
            content_type = "text/html" if ext == ".html" else "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            supabase.storage.from_("templates").upload(
                path=lor_path,
                file=lor_bytes,
                file_options={"content-type": content_type, "upsert": "true"}
            )
        
        if experience_template:
            exp_bytes = await experience_template.read()
            ext = os.path.splitext(experience_template.filename)[1].lower()
            exp_path = f"templates/{batch_id}/experience{ext}"
            content_type = "text/html" if ext == ".html" else "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            supabase.storage.from_("templates").upload(
                path=exp_path,
                file=exp_bytes,
                file_options={"content-type": content_type, "upsert": "true"}
            )

        if internship_template:
            int_bytes = await internship_template.read()
            ext = os.path.splitext(internship_template.filename)[1].lower()
            int_path = f"templates/{batch_id}/internship{ext}"
            content_type = "text/html" if ext == ".html" else "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            supabase.storage.from_("templates").upload(
                path=int_path,
                file=int_bytes,
                file_options={"content-type": content_type, "upsert": "true"}
            )

        batch_record = {
            "id": batch_id,
            "month": month,
            "issue_date": issue_date,
            "lor_template_path": lor_path,
            "experience_template_path": exp_path,
            "internship_template_path": int_path
        }
        supabase.table("batches").upsert(batch_record).execute()

        return {"status": "success", "batch_id": batch_id}

    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to create batch: {str(e)}"
        )






@app.post("/api/generate")
async def generate_certificates(
    background_tasks: BackgroundTasks,
    excel_file: UploadFile = File(...),
    batch_id: Optional[str] = Form(None),
    template_id: Optional[str] = Form(None)
):
    """
    Processes the details from Excel.
    For each row, inserts records into the 'interns' and 'certificates' tables
    based on their LOR, Experience, and Internship selections.
    Emails certificates to interns.
    """
    target_batch_id = (batch_id or template_id or "").strip()
    if not target_batch_id:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Missing batch_id or template_id parameter."
        )

    if not supabase:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Supabase client is not configured."
        )

    # Validate Excel format
    if not excel_file.filename.lower().endswith((".xlsx", ".xls")):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Please upload a valid Excel file (.xlsx or .xls)."
        )

    # Read and parse Excel file
    try:
        excel_bytes = await excel_file.read()
        df = pd.read_excel(io.BytesIO(excel_bytes))
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to parse Excel file: {e}"
        )

    # Clean and validate Excel columns
    df.columns = [str(col).strip().lower() for col in df.columns]
    col_mapping = {}
    for col in df.columns:
        if "name" in col:
            col_mapping["name"] = col
        elif "email" in col:
            col_mapping["email"] = col
        elif "college" in col or "university" in col or "institution" in col:
            col_mapping["college"] = col
        elif "year" in col or "class" in col:
            col_mapping["year"] = col
        elif "department" in col or "dept" in col or "branch" in col:
            col_mapping["department"] = col
        elif "role" in col or "domain" in col:
            col_mapping["role"] = col
        elif "project" in col or "proj" in col or "area" in col:
            col_mapping["project"] = col
        elif "month" in col or "batch" in col:
            col_mapping["month"] = col
        elif "date" in col or "dt" in col:
            col_mapping["date"] = col
        
        # YES/NO certificate selections
        elif "recomandation" in col or "recommendation" in col or "lor" in col:
            col_mapping["opt_lor"] = col
        elif "experience" in col:
            col_mapping["opt_experience"] = col
        elif "internship certificate" in col or "internship_certificate" in col:
            col_mapping["opt_internship"] = col

    # Fallback to standard column mapping check
    missing_cols = []
    if "name" not in col_mapping:
        missing_cols.append("Name")
    if "email" not in col_mapping:
        missing_cols.append("Email")
    if "college" not in col_mapping:
        missing_cols.append("College")
    if "year" not in col_mapping:
        missing_cols.append("Year")
    if "department" not in col_mapping:
        missing_cols.append("Department")

    if missing_cols:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Excel file is missing required columns: {', '.join(missing_cols)}"
        )

    # Process row by row
    rows_results = []

    for index, row in df.iterrows():
        try:
            def _cell_str(r, key):
                if key not in col_mapping:
                    return ""
                v = r.get(col_mapping[key], "")
                if pd.isna(v):
                    return ""
                if isinstance(v, (pd.Timestamp,)):
                    return str(v.date())
                return str(v).strip()

            def _cell_yes_no(r, key):
                val = _cell_str(r, key).upper()
                return val == "YES"

            name_val = _cell_str(row, "name")
            email_val = _cell_str(row, "email")
            college_val = _cell_str(row, "college")
            year_val = _cell_str(row, "year")
            department_val = _cell_str(row, "department")
            role_val = _cell_str(row, "role")
            project_val = _cell_str(row, "project")
            month_val = _cell_str(row, "month")
            date_val = _cell_str(row, "date")

            opt_lor = _cell_yes_no(row, "opt_lor")
            opt_exp = _cell_yes_no(row, "opt_experience")
            opt_int = _cell_yes_no(row, "opt_internship")

            if not name_val or name_val.lower() == "nan" or not email_val:
                continue

            # Parse Issue Date
            issue_date_val = None
            if date_val:
                try:
                    clean_dt = date_val.strip()
                    parsed_dt = pd.to_datetime(clean_dt, dayfirst=True, errors='raise')
                    issue_date_val = parsed_dt.date().isoformat()
                except Exception:
                    try:
                        from dateutil import parser
                        parsed_dt = parser.parse(clean_dt, dayfirst=True)
                        issue_date_val = parsed_dt.date().isoformat()
                    except Exception:
                        issue_date_val = None

            # 1. Insert/Save Intern details
            intern_data = {
                "email": email_val,
                "name": name_val,
                "college": college_val,
                "year": year_val,
                "department": department_val,
                "role": role_val,
                "project": project_val,
                "month": month_val,
                "date": date_val,
                "batch_id": target_batch_id
            }

            intern_res = supabase.table("interns").insert(intern_data).execute()
            if not intern_res.data:
                raise ValueError("Failed to create intern database record.")
            intern_id = intern_res.data[0]["id"]

            # 2. Build Certificates based on YES/NO choices
            certificates_to_create = []
            if opt_lor:
                certificates_to_create.append(("lor", "Letter of Recommendation"))
            if opt_exp:
                certificates_to_create.append(("experience", "Experience Letter"))
            if opt_int:
                certificates_to_create.append(("internship", "Internship Certificate"))

            email_certificates = []

            for cert_type, label in certificates_to_create:
                cert_code = str(uuid.uuid4())
                verify_url = build_verify_url(cert_code)
                dynamic_pdf_url = f"{BACKEND_BASE_URL}/certificate/{cert_code}.pdf"

                cert_record = {
                    "cert_code": cert_code,
                    "intern_id": intern_id,
                    "cert_type": cert_type,
                    "status": "active",
                    "pdf_url": dynamic_pdf_url,
                    "name": name_val,
                    "college": college_val,
                    "batch": year_val,
                    "department": department_val,
                    "role": role_val,
                    "project": project_val,
                    "month": month_val,
                    "issue_date": date_val
                }

                supabase.table("certificates").insert(cert_record).execute()

                email_certificates.append({
                    "type": cert_type,
                    "label": label,
                    "cert_code": cert_code,
                    "pdf_url": dynamic_pdf_url,
                    "verify_url": verify_url
                })

                # Append to rows results for frontend UI display
                rows_results.append({
                    "name": name_val,
                    "college": college_val,
                    "department": department_val,
                    "month": f"{label} ({month_val or year_val})",
                    "cert_code": cert_code,
                    "verify_url": verify_url,
                    "pdf_url": dynamic_pdf_url,
                    "status": "active",
                    "intern_id": str(intern_id),
                    "email": email_val,
                    "email_status": "pending"
                })

        except Exception as row_error:
            print(f"Row {index} failed: {row_error}")
            rows_results.append({
                "name": str(row.get(col_mapping.get("name", ""), "Unknown")),
                "college": str(row.get(col_mapping.get("college", ""), "Unknown")),
                "department": str(row.get(col_mapping.get("department", ""), "Unknown")),
                "month": "Error",
                "cert_code": None,
                "pdf_url": None,
                "status": "error",
                "error": str(row_error)
            })

    return {
        "excel_download_url": "",
        "rows": rows_results
    }


@app.get("/api/certificates/{cert_code}")
def lookup_certificate(cert_code: str):
    """
    Looks up details of a certificate from database (supports case-insensitive matching).
    """
    if not supabase:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Supabase client is not configured."
        )

    try:
        # Try exact match first
        res = supabase.table("certificates").select("*").eq("cert_code", cert_code).execute()
        if not res.data:
            # Fallback to uppercase
            res = supabase.table("certificates").select("*").eq("cert_code", cert_code.upper()).execute()
        if not res.data:
            # Fallback to lowercase
            res = supabase.table("certificates").select("*").eq("cert_code", cert_code.lower()).execute()

        if not res.data:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"Certificate with code {cert_code} not found."
            )
        
        cert_data = dict(res.data[0])
        intern_id = cert_data.get("intern_id")
        if intern_id:
            intern_res = supabase.table("interns").select("*").eq("id", intern_id).execute()
            if intern_res.data:
                intern = intern_res.data[0]
                cert_data["name"] = intern.get("name")
                cert_data["college"] = intern.get("college")
                cert_data["batch"] = intern.get("year")
                cert_data["department"] = intern.get("department")
                cert_data["role"] = intern.get("role")
                cert_data["project"] = intern.get("project")
                cert_data["month"] = intern.get("month")
                cert_data["issue_date"] = intern.get("date") or cert_data.get("issue_date")
        
        return cert_data

    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Database lookup failed: {str(e)}"
        )


async def get_pdf_bytes_for_certificate(cert_code: str) -> bytes:
    # 0. Check PDF cache first
    cached_pdf = pdf_cache.get(cert_code)
    if cached_pdf:
        return cached_pdf

    if not supabase:
        raise ValueError("Supabase client is not configured.")

    # 1. Fetch certificate record
    res = supabase.table("certificates").select("*").eq("cert_code", cert_code).execute()
    if not res.data:
        res = supabase.table("certificates").select("*").eq("cert_code", cert_code.upper()).execute()
    if not res.data:
        res = supabase.table("certificates").select("*").eq("cert_code", cert_code.lower()).execute()

    if not res.data:
        raise ValueError(f"Certificate with code {cert_code} not found.")
    
    cert_data = res.data[0]
    intern_id = cert_data.get("intern_id")
    cert_type = cert_data.get("cert_type") or "internship"
    batch_id = cert_data.get("template_id") or cert_data.get("batch_id")

    intern_name = cert_data.get("name")
    college = cert_data.get("college")
    year = cert_data.get("batch")
    department = cert_data.get("department")
    role = cert_data.get("role")
    project = cert_data.get("project")
    month = cert_data.get("month")
    date_val = cert_data.get("issue_date")

    if intern_id:
        intern_res = supabase.table("interns").select("*").eq("id", intern_id).execute()
        if intern_res.data:
            intern = intern_res.data[0]
            intern_name = intern.get("name") or intern_name
            college = intern.get("college") or college
            year = intern.get("year") or year
            department = intern.get("department") or department
            role = intern.get("role") or role
            project = intern.get("project") or project
            month = intern.get("month") or month
            date_val = intern.get("date") or date_val
            batch_id = intern.get("batch_id") or batch_id

    if not batch_id:
        if cert_type == "lor":
            batch_id = "slot-2"
        elif cert_type == "experience":
            batch_id = "slot-3"
        else:
            batch_id = "slot-1"

    replacements = {
        "<<NAME>>": intern_name or "",
        "<<INSTITUTION>>": college or "",
        "<<COLLEGE>>": college or "",
        "<<YEAR>>": year or "",
        "<<DEPARTMENT>>": department or "",
        "<<DOMAIN>>": role or "",
        "<<ROLE>>": role or "",
        "<<PROJECT>>": project or "",
        "<<INTERNSHIP & LIVE PROJECT AREA>>": project or "",
        "<<BATCH>>": month or "",
        "<<BATCH >>": month or "",
        "<<DATE>>": str(date_val) if date_val else "",
        "<<DT>>": str(date_val) if date_val else ""
    }

    qr_url = build_verify_url(cert_code)
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_H,
        box_size=12,
        border=2
    )
    qr.add_data(qr_url)
    qr.make(fit=True)
    from qrcode.image.pil import PilImage
    qr_img = qr.make_image(image_factory=PilImage, fill_color="black", back_color="white")
    qr_io = io.BytesIO()
    try:
        qr_img.save(qr_io, format="PNG")
    except TypeError:
        qr_io.seek(0)
        qr_io.truncate(0)
        qr_img.save(qr_io)
    qr_bytes = qr_io.getvalue()


    # 3. Resolve the .pptx template file path and download it.
    #    We go directly to the .pptx — NEVER use the overlay pipeline,
    #    because the overlay only renders dynamic text and drops all static
    #    elements (footer, logo, decorative lines, signature image).
    template_bytes = None
    ext = ".pptx"

    # Priority 1: named file e.g. templates/BATCH_TEST1/internship.pptx
    if batch_id and not batch_id.startswith("slot-"):
        import re as _re
        batch_id_stripped = batch_id.strip()
        batch_id_normalized = _re.sub(r'\s+', '_', batch_id_stripped.lower())
        cert_type_file = f"{cert_type}.pptx"

        for folder in [batch_id_stripped, batch_id_normalized]:
            if not folder:
                continue
            try:
                candidate = f"templates/{folder}/{cert_type_file}"
                template_bytes = supabase.storage.from_("templates").download(candidate)
                print(f"[pdf] Using named PPTX: {candidate}")
                break
            except Exception:
                pass

    # Priority 2: slot-based path e.g. templates/slot-1.pptx
    if not template_bytes:
        slot_path = f"templates/{batch_id}.pptx"
        try:
            template_bytes = supabase.storage.from_("templates").download(slot_path)
            print(f"[pdf] Using slot PPTX: {slot_path}")
        except Exception:
            pass

    # Priority 3: batches table lookup
    if not template_bytes:
        try:
            batch_res = supabase.table("batches").select("*").eq("id", batch_id).execute()
            if batch_res.data:
                batch = batch_res.data[0]
                if cert_type == "lor":
                    template_path = batch.get("lor_template_path")
                elif cert_type == "experience":
                    template_path = batch.get("experience_template_path")
                else:
                    template_path = batch.get("internship_template_path")
                if template_path:
                    ext = os.path.splitext(template_path)[1].lower() or ".pptx"
                    template_bytes = supabase.storage.from_("templates").download(template_path)
                    print(f"[pdf] Using batches-table path: {template_path}")
        except Exception:
            pass

    if not template_bytes:
        raise ValueError(f"Template PPTX not found for batch '{batch_id}' / cert_type '{cert_type}'. "
                         f"Upload templates/{{batch_id}}/{{cert_type}}.pptx to storage.")

    # 4. Use HTML/CSS WeasyPrint layout engine conversion (Technique 1).
    #    This preserves every font, color, image, line, footer — EVERYTHING
    #    from the original designer template while providing absolute reflow/wrapping.
    if ext == ".pptx":
        pdf_bytes = await generate_pdf_bytes_with_weasyprint_async(batch_id, cert_type, template_bytes, replacements, qr_bytes)
    else:
        pdf_bytes, _ = generate_certificate_from_html_bytes(template_bytes, replacements, qr_bytes)

    pdf_cache.set(cert_code, pdf_bytes)
    return pdf_bytes



@app.get("/certificate/{cert_code}.pdf")
@app.get("/api/certificates/{cert_code}/pdf")
async def get_dynamic_pdf(cert_code: str):
    """
    Generates and returns the PDF for a specific certificate dynamically.
    Utilizes template caching, transient PDF caching, and serialized LibreOffice conversion.
    """
    try:
        safe_name = "certificate"
        try:
            res = supabase.table("certificates").select("name").eq("cert_code", cert_code).execute()
            if res.data:
                safe_name = re.sub(r'[\\/*?:"<>|]', "_", str(res.data[0].get("name") or "certificate")).strip()
        except:
            pass

        pdf_bytes = await get_pdf_bytes_for_certificate(cert_code)
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"inline; filename={safe_name}_{cert_code}.pdf"
            }
        )
    except ValueError as ve:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail=str(ve))
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to generate dynamic certificate PDF: {str(e)}"
        )




@app.post("/api/interns/{intern_id}/send-email")
async def manual_send_intern_email(intern_id: str):
    """
    Manually triggers email dispatch for a specific intern.
    Loads intern details and all associated certificates, then sends SMTP email.
    """
    if not supabase:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Supabase client is not configured."
        )

    try:
        # 1. Fetch Intern details
        intern_res = supabase.table("interns").select("*").eq("id", intern_id).execute()
        if not intern_res.data:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"Intern with ID {intern_id} not found."
            )
        intern_data = intern_res.data[0]
        email_val = intern_data.get("email")
        name_val = intern_data.get("name")
        month_val = intern_data.get("month")
        year_val = intern_data.get("year")

        if not email_val:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Intern does not have a configured email address."
            )

        # 2. Fetch associated certificates
        certs_res = supabase.table("certificates").select("*").eq("intern_id", intern_id).execute()
        if not certs_res.data:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"No certificates found associated with intern ID {intern_id}."
            )

        email_certificates = []
        for c in certs_res.data:
            cert_type = c.get("cert_type") or "internship"
            label = "Certificate"
            if cert_type == "lor":
                label = "Letter of Recommendation"
            elif cert_type == "experience":
                label = "Experience Letter"
            elif cert_type == "internship":
                label = "Internship Certificate"

            email_certificates.append({
                "type": cert_type,
                "label": label,
                "cert_code": c.get("cert_code"),
                "pdf_url": c.get("pdf_url")
            })

        # 3. Call email notification
        success = await send_email_notification(
            to_email=email_val,
            intern_name=name_val,
            batch_title=month_val or year_val,
            certificates=email_certificates,
            intern_id=intern_id
        )

        # Update email_status in the DB if the column exists (defensive try-except)
        next_status = "sent" if success else "failed"
        try:
            supabase.table("interns").update({"email_status": next_status}).eq("id", intern_id).execute()
        except Exception as update_err:
            print(f"Could not update email_status in database: {update_err}")

        if not success:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="SMTP dispatch failed. Check backend credentials and SMTP server config."
            )

        return {"status": "success", "message": f"Email successfully dispatched to {email_val}."}

    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to dispatch email: {str(e)}"
        )


@app.post("/api/certificates/{cert_code}/revoke")
def revoke_certificate(cert_code: str, reason: Optional[str] = None):
    """
    Admin-only command to revoke a certificate.
    Sets status = 'revoked'.
    """
    if not supabase:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Supabase client is not configured."
        )

    try:
        update_data = {
            "status": "revoked",
            "revoked_at": datetime.utcnow().isoformat()
        }
        if reason:
            update_data["revoke_reason"] = reason

        # Find case-insensitive target code
        target_code = cert_code
        check_res = supabase.table("certificates").select("cert_code").eq("cert_code", cert_code).execute()
        if not check_res.data:
            check_res_upper = supabase.table("certificates").select("cert_code").eq("cert_code", cert_code.upper()).execute()
            if check_res_upper.data:
                target_code = cert_code.upper()
            else:
                check_res_lower = supabase.table("certificates").select("cert_code").eq("cert_code", cert_code.lower()).execute()
                if check_res_lower.data:
                    target_code = cert_code.lower()

        res = supabase.table("certificates").update(update_data).eq("cert_code", target_code).execute()
        if not res.data:
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND,
                detail=f"Certificate with code {cert_code} not found."
            )
        return {"status": "success", "message": f"Certificate {target_code} revoked successfully."}
    except HTTPException as he:
        raise he
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Database update failed: {str(e)}"
        )


if __name__ == "__main__":
    import uvicorn

    port = int(os.getenv("PORT", 5000))
    # Bind to 0.0.0.0 so the dev server is reachable from the host machine
    # Use import string so automatic reload works when running `python app.py`
    uvicorn.run("app:app", host="0.0.0.0", port=port, reload=True)
