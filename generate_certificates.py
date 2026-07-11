import os
import re
import uuid
import argparse
import io
import shutil
import subprocess
import pandas as pd
import qrcode
from pptx import Presentation
from pptx.util import Inches

# win32com has been removed as this tool is adapted for headless Linux

# =====================================================
# HELPERS
# =====================================================

def safe_filename(name):
    """Sanitize filenames for filesystem safety."""
    return re.sub(r'[\\/*?:"<>|]', "_", str(name)).strip()

def generate_qr_code(payload, out_path):
    """Generate a high-quality, high-resolution scannable QR code."""
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_H,  # High error correction
        box_size=12,  # Ensures resolution is at least 300x300px
        border=2      # Non-zero quiet zone/border
    )
    qr.add_data(payload)
    qr.make(fit=True)
    
    from qrcode.image.pil import PilImage
    img = qr.make_image(image_factory=PilImage, fill_color="black", back_color="white")
    img.save(out_path)

def export_pptx_to_pdf(pptx_path: str, output_dir: str) -> str:
    """
    Converts PPTX to PDF using headless LibreOffice.
    Uses an isolated user profile per call to avoid lock-file
    collisions under concurrent requests.
    """
    os.makedirs(output_dir, exist_ok=True)

    # Isolated profile dir per invocation — prevents concurrent
    # soffice processes from fighting over the same lock file
    profile_dir = f"/tmp/lo_profile_{uuid.uuid4().hex}"

    cmd = [
        "soffice",
        "--headless",
        "--norestore",
        "--nofirststartwizard",
        "--nologo",
        "--nodefault",
        "--invisible",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        f"-env:UserInstallation=file://{profile_dir}",
        pptx_path,
    ]

    # Explicitly set HOME environment variable to /tmp to prevent LibreOffice
    # from failing when trying to write to non-writable Nixpacks/Dokploy directories.
    env = os.environ.copy()
    env["HOME"] = "/tmp"

    try:
        subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=60,  # hard timeout — do not let a stuck soffice process hang
            check=True,
            env=env,
        )
    except subprocess.TimeoutExpired:
        raise RuntimeError(
            f"LibreOffice conversion timed out for {pptx_path}"
        )
    except subprocess.CalledProcessError as e:
        raise RuntimeError(
            f"LibreOffice conversion failed with exit code {e.returncode}.\nStdout: {e.stdout}\nStderr: {e.stderr}"
        )
    finally:
        # Always clean up the temp profile directory, even on failure
        shutil.rmtree(profile_dir, ignore_errors=True)

    expected_pdf = os.path.join(
        output_dir,
        os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf",
    )

    if not os.path.exists(expected_pdf):
        raise RuntimeError(
            f"Conversion reported success but output PDF not found: "
            f"{expected_pdf}"
        )

    return expected_pdf

def convert_pptx_to_pdf(pptx_path, pdf_path):
    """Convert PPTX to PDF using headless LibreOffice."""
    try:
        abs_pptx = os.path.abspath(pptx_path)
        abs_pdf = os.path.abspath(pdf_path)
        out_dir = os.path.dirname(abs_pdf)
        
        expected_pdf = export_pptx_to_pdf(abs_pptx, out_dir)
        
        if os.path.abspath(expected_pdf) != abs_pdf:
            if os.path.exists(abs_pdf):
                os.remove(abs_pdf)
            os.rename(expected_pdf, abs_pdf)
        return True
    except Exception as e:
        print(f"  --> LibreOffice conversion failed: {e}")
        return False

def defragment_paragraph(paragraph, placeholders):
    """
    Defragments python-pptx runs that have split placeholders.
    Converts split runs into a single run so they can be replaced accurately.
    """
    for key in placeholders:
        keys_to_check = [key, key.replace("<<", "«").replace(">>", "»")]
        for k_check in keys_to_check:
            search_start = 0
            while True:
                p_text = "".join(run.text for run in paragraph.runs)
                start_idx = p_text.find(k_check, search_start)
                if start_idx == -1:
                    break

                # Map each character position to its run index
                run_map = []
                for run_idx, run in enumerate(paragraph.runs):
                    run_map.extend([run_idx] * len(run.text))

                end_idx = start_idx + len(k_check) - 1
                if end_idx >= len(run_map):
                    break

                # Get run indices that cover the placeholder
                run_indices = run_map[start_idx : end_idx + 1]
                first_run_idx = run_indices[0]
                last_run_idx = run_indices[-1]

                if first_run_idx == last_run_idx:
                    search_start = start_idx + len(k_check)
                    continue

                # Merge split runs
                combined_text = "".join(paragraph.runs[r_i].text for r_i in range(first_run_idx, last_run_idx + 1))
                paragraph.runs[first_run_idx].text = combined_text
                for r_i in range(first_run_idx + 1, last_run_idx + 1):
                    paragraph.runs[r_i].text = ""

                search_start = start_idx + len(k_check)

def text_frame_has_placeholders(text_frame, replacements) -> bool:
    """
    Returns True if the text frame contains any of the keys in replacements,
    or general QR placeholders.
    """
    text = text_frame.text
    for key in replacements:
        if key in text:
            return True
        guill_key = key.replace("<<", "«").replace(">>", "»")
        if guill_key in text:
            return True
    
    for qr_key in ("<<QR>>", "<<QR_CODE>>", "<<QRCODE>>", "«QR»", "«QR_CODE»", "«QRCODE»"):
        if qr_key in text:
            return True
            
    return False

def detect_certificate_title(slide) -> str:
    """
    Detects the certificate type based on the text contents of the slide shapes.
    Returns: 'Internship Certificate', 'Letter Of Recomandation', 'Experience Letter', or 'Certificate'.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if "recommendation" in text or "recomandation" in text:
                return "Letter Of Recomandation"
            elif "experience certificate" in text or "experience letter" in text:
                return "Experience Letter"
            elif "internship" in text:
                return "Internship Certificate"
    
    # Fallback checks
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            if "certificate" in text:
                return "Certificate"
                
    return "Certificate"

# =====================================================
# MAIN PIPELINE
# =====================================================

def main():
    parser = argparse.ArgumentParser(description="PPTX to PDF Certificate Generator Batch Tool")
    parser.add_argument(
        "--template", 
        default="C:/Users/ROHITH P/Downloads/TEMPLATE.pptx",
        help="Path to PPTX certificate template slide"
    )
    parser.add_argument(
        "--data", 
        default="C:/Users/ROHITH P/Desktop/TEST 3.xlsx",
        help="Path to Excel spreadsheet rows"
    )
    parser.add_argument(
        "--outdir", 
        default="output",
        help="Directory to save generated certificates"
    )
    args = parser.parse_args()

    # Validate inputs
    if not os.path.exists(args.template):
        # Check current folder if default download path doesn't exist
        if os.path.exists("TEMPLATE.pptx"):
            args.template = "TEMPLATE.pptx"
        else:
            raise FileNotFoundError(f"Template PPTX not found: {args.template}")

    if not os.path.exists(args.data):
        if os.path.exists("TEST 3.xlsx"):
            args.data = "TEST 3.xlsx"
        else:
            raise FileNotFoundError(f"Excel data sheet not found: {args.data}")

    os.makedirs(args.outdir, exist_ok=True)
    temp_dir = os.path.join(args.outdir, "temp_assets")
    os.makedirs(temp_dir, exist_ok=True)

    print("=" * 60)
    print("Starting PPTX Certificate Generator...")
    print(f"Template   : {args.template}")
    print(f"Data Sheet : {args.data}")
    print(f"Out Dir    : {args.outdir}")
    print("=" * 60)

    # 1. Load Excel and trim headers/cell values
    df = pd.read_excel(args.data)
    df.columns = [str(col).strip() for col in df.columns]
    
    # Map headers robustly
    col_mapping = {}
    for col in df.columns:
        col_lower = col.lower()
        if "name" in col_lower:
            col_mapping["name"] = col
        elif "college" in col_lower or "university" in col_lower or "institution" in col_lower:
            col_mapping["college"] = col
        elif "year" in col_lower or "class" in col_lower:
            col_mapping["year"] = col
        elif "department" in col_lower or "dept" in col_lower or "branch" in col_lower:
            col_mapping["department"] = col
        elif "role" in col_lower or "domain" in col_lower:
            col_mapping["role"] = col
        elif "project" in col_lower or "proj" in col_lower or "area" in col_lower:
            col_mapping["project"] = col
        elif "month" in col_lower or "batch" in col_lower:
            col_mapping["month"] = col
        elif "date" in col_lower or "dt" in col_lower:
            col_mapping["date"] = col

    # Check for name column
    if "name" not in col_mapping:
        raise ValueError("Excel file is missing a 'NAME' column.")

    def _get_row_val(row_data, key):
        if key not in col_mapping:
            return ""
        val = row_data.get(col_mapping[key], "")
        if pd.isna(val):
            return ""
        return str(val).strip()

    # 2. Process rows
    report = []
    total_rows = len(df)
    
    for index, row in df.iterrows():
        name_val = _get_row_val(row, "name")
        if not name_val or name_val.lower() == "nan":
            continue
            
        college_val = _get_row_val(row, "college")
        year_val = _get_row_val(row, "year")
        department_val = _get_row_val(row, "department")
        role_val = _get_row_val(row, "role")
        project_val = _get_row_val(row, "project")
        month_val = _get_row_val(row, "month")
        date_val = _get_row_val(row, "date")
        
        # Clean date format if parsed
        if date_val:
            try:
                parsed_dt = pd.to_datetime(date_val)
                date_val = parsed_dt.strftime("%d.%m.%Y")
            except Exception:
                pass

        cert_code = str(uuid.uuid4())
        print(f"[{index + 1}/{total_rows}] Processing: {name_val}")
        
        # Copy template presentation
        out_pptx_name = f"{safe_filename(name_val)}_{index}.pptx"
        out_pptx_path = os.path.join(args.outdir, out_pptx_name)
        shutil.copy2(args.template, out_pptx_path)
        
        # Setup replacements dictionary
        replacements = {
            "<<NAME>>": name_val,
            "<<INSTITUTION>>": college_val,
            "<<COLLEGE>>": college_val,
            "<<YEAR>>": year_val,
            "<<DEPARTMENT>>": department_val,
            "<<DOMAIN>>": role_val,
            "<<ROLE>>": role_val,
            "<<PROJECT>>": project_val,
            "<<INTERNSHIP & LIVE PROJECT AREA>>": project_val,
            "<<BATCH>>": month_val,
            "<<BATCH >>": month_val,
            "<<DATE>>": date_val,
            "<<DT>>": date_val
        }
        
        # Generate temporary QR code PNG
        qr_path = os.path.join(temp_dir, f"qr_{cert_code}.png")
        verify_url = f"https://coptercode.co.in/verify?id={cert_code}"
        generate_qr_code(verify_url, qr_path)
        
        # Open PPTX in python-pptx
        prs = Presentation(out_pptx_path)
        slide = prs.slides[0]
        
        qr_placed = False
        
        # Fields that should WRAP to next line instead of shrinking font
        WRAP_FIELDS = {"<<INTERNSHIP & LIVE PROJECT AREA>>", "<<PROJECT>>", "<<DOMAIN>>", "<<ROLE>>"}

        # Find QR shape and replace text
        for shape in list(slide.shapes):
            if shape.has_text_frame:
                text_frame = shape.text_frame
                combined_text = text_frame.text.strip()

                # Fix for "COPTERCODE" wrapping "E" to the next line in the header textbox
                if ("COPTERCODE" in combined_text or "COPTERCOD" in combined_text) and len(combined_text) < 15:
                    text_frame.word_wrap = False
                    shape.width = shape.width + Inches(0.5)

                # Skip text boxes that don't contain any placeholders at all to preserve formatting/alignment
                if not text_frame_has_placeholders(text_frame, replacements):
                    continue

                # Enable word wrap so long values flow to next line naturally
                text_frame.word_wrap = True

                # Defragment paragraph runs to resolve split placeholders
                placeholders_to_defrag = list(replacements.keys()) + ["<<QR>>", "<<QR_CODE>>", "<<QRCODE>>"]
                for paragraph in text_frame.paragraphs:
                    defragment_paragraph(paragraph, placeholders_to_defrag)

                # Check for QR placeholder
                combined_text = text_frame.text.strip()
                if "<<QR>>" in combined_text or "«QR»" in combined_text:
                    # Make the QR code larger (1.15 in x 1.15 in)
                    qr_w = Inches(1.15)
                    qr_h = Inches(1.15)
                    
                    # Center the QR image relative to the original <<QR>> placeholder textbox
                    qr_left = int(shape.left + (shape.width - qr_w) // 2)
                    qr_top = int(shape.top + (shape.height - qr_h) // 2)

                    # Remove the original text shape
                    sp = shape._element
                    sp.getparent().remove(sp)

                    # Add picture (Aspect-ratio locked, 1.15 in x 1.15 in)
                    slide.shapes.add_picture(
                        qr_path,
                        qr_left,
                        qr_top,
                        width=qr_w,
                        height=qr_h
                    )
                    qr_placed = True
                    continue

                # Replace tokens in-place in text runs (preserves formatting)
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, val in replacements.items():
                            has_match = False
                            actual_key = None

                            if key in run.text:
                                has_match = True
                                actual_key = key
                            else:
                                guill_key = key.replace("<<", "«").replace(">>", "»")
                                if guill_key in run.text:
                                    has_match = True
                                    actual_key = guill_key

                            if has_match:
                                # Standardize to Arial to prevent system font substitutions
                                run.font.name = "Arial"

                                # For body/area fields: keep font size and let text wrap
                                # For short label fields (name, institution): allow mild shrinking
                                is_wrap_field = key in WRAP_FIELDS
                                if not is_wrap_field and run.font.size and val:
                                    length = len(val)
                                    limit = 20 if "NAME" in key else 30
                                    if length > limit:
                                        scale = limit / length
                                        scale = max(0.75, min(1.0, scale))
                                        run.font.size = int(run.font.size * scale)

                                run.text = run.text.replace(actual_key, val)

                                # Clean up visual hacks (multiple consecutive spaces) and fix missing comma spaces
                                if run.text:
                                    run.text = re.sub(r',([a-zA-Z])', r', \1', run.text)
                                    run.text = re.sub(r'\s{2,}', ' ', run.text)
                                
        # Detect certificate title from slide text
        cert_title = detect_certificate_title(slide)

        # Save presentation
        prs.save(out_pptx_path)
        
        # Convert to PDF
        out_pdf_name = f"{safe_filename(name_val)}_({cert_title}).pdf"
        out_pdf_path = os.path.join(args.outdir, out_pdf_name)
        
        success = convert_pptx_to_pdf(out_pptx_path, out_pdf_path)
        
        # Optionally delete intermediate PPTX to clean up output folder
        if os.path.exists(out_pptx_path):
            os.remove(out_pptx_path)
            
        report.append({
            "NAME": name_val,
            "CERT_CODE": cert_code,
            "VERIFY_URL": verify_url,
            "STATUS": "SUCCESS" if success else "FAILED",
            "PDF": os.path.abspath(out_pdf_path) if success else ""
        })
        
    # Clean up temp assets folder
    shutil.rmtree(temp_dir, ignore_errors=True)
    
    # Save compilation report
    report_df = pd.DataFrame(report)
    report_df.to_csv(os.path.join(args.outdir, "report.csv"), index=False)
    
    print("\n" + "=" * 60)
    print("Batch Run Complete!")
    print(f"Generated {len(report_df[report_df['STATUS'] == 'SUCCESS'])}/{len(report_df)} certificates.")
    print(f"Output saved in: {os.path.abspath(args.outdir)}")
    print("=" * 60)

if __name__ == '__main__':
    main()
